#!/usr/bin/env python3
"""
CData Final Round — Principal Developer Marketing Manager
==========================================================
30-minute presentation + 10-minute Q&A

Run:
    cd presentation
    ../.venv/bin/python3 build_slides.py
    open cdata-interview.pptx

Slide order:
  HOOK (5 slides)
    1.  Title
    2.  Before the strategy
    3.  Live demo
    4.  Demo result
    5.  Built it. Broke it.

  PART 1: The AI Developer Opportunity (10 slides)
    6.  Developer who needs something specific
    7.  Cross-source pain
    8.  Enterprise data (Salesforce, SAP, Workday, NetSuite)
    9.  CData position — two decades of enterprise connectors
   10.  Not reached yet
   11.  Three audiences intro
   12.  Audience — LLM agents
   13.  Audience — AI developers
   14.  Audience — Enterprise buyers
   15.  Each needs a different approach

  PART 2: The Audit (16 slides)
   16.  What's actually strong
   17.  Complete issue table (47 issues, 6 categories)
   18.  Not walking through all 47
   19.  Finding 1 intro — asked Claude Code
   20.  Screenshot — Composio recommended
   21.  5-query table — CData in 0/5
   22.  LLMs are where this developer starts
   23.  Finding 2 — no agent-specific path
   24.  Same seven steps
   25.  CTAs, not docs
   26.  Finding 3 — Connect AI effectively orphaned
   27.  No path to docs
   28.  Wrong place
   29.  The pattern
   30.  Two buckets
   31.  PMM fixes vs product conversations

  PART 3: The GTM Plan (10 slides)
   32.  90 days
   33.  Days 1-30 intro
   34.  Three parallel workstreams
   35.  Immediate fixes that don't wait
   36.  AI developer entry point
   37.  Days 31-60
   38.  Days 61-90
   39.  Goals framework
   40.  PMM / DevRel split
   41.  What developer content actually works

  CLOSE (3 slides)
   42.  What I'm leaving behind
   43.  Close statement
   44.  Questions
"""

import sys, os
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from pptx_utils import *
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

# ── CData brand palette ───────────────────────────────────────────────────────
CDATA_YELLOW = RGBColor(0xFF, 0xE5, 0x00)
CDATA_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
CDATA_MID    = RGBColor(0x4A, 0x4A, 0x4A)
CDATA_GRAY   = RGBColor(0x80, 0x80, 0x80)
CDATA_CARD   = RGBColor(0xEC, 0xED, 0xE8)
CDATA_BG     = RGBColor(0xF5, 0xF5, 0xF0)
CDATA_DIVIDER= RGBColor(0xD0, 0xD0, 0xC8)

# Shadow pptx_utils globals so layout helpers pick up CData colors
GREEN    = CDATA_YELLOW
BG_LIME  = CDATA_YELLOW
BG       = CDATA_BG
DARK     = CDATA_DARK
LGRAY    = CDATA_MID
GRAY     = CDATA_GRAY
CARD     = CDATA_CARD
DIVIDER  = CDATA_DIVIDER
TEAL     = CDATA_DARK


# ── Logo generation ───────────────────────────────────────────────────────────

def _generate_cdata_logos():
    from PIL import Image, ImageDraw, ImageFont as _IFont
    _FONT = '/System/Library/Fonts/HelveticaNeue.ttc'

    def _make(variant):
        img = Image.new('RGBA', (600, 100), (0, 0, 0, 0))
        d   = ImageDraw.Draw(img)
        if variant == 'light':
            mark_c, text_c = (0xFF, 0xE5, 0x00, 255), (0xFF, 0xFF, 0xFF, 255)
        elif variant == 'yellow':
            mark_c, text_c = (0x1A, 0x1A, 0x1A, 255), (0x1A, 0x1A, 0x1A, 255)
        else:
            mark_c, text_c = (0xFF, 0xE5, 0x00, 255), (0x1A, 0x1A, 0x1A, 255)
        try:
            font = _IFont.truetype(_FONT, 52, index=1)
        except Exception:
            font = _IFont.load_default()
        d.rectangle([4, 18, 48, 74], fill=mark_c)
        d.rectangle([22, 30, 48, 62], fill=(0, 0, 0, 0))
        d.text((62, 16), 'CData', font=font, fill=text_c)
        bbox = img.getbbox()
        return img.crop(bbox) if bbox else img

    for v in ('dark', 'light', 'yellow'):
        _make(v).save(os.path.join(HERE, f'cdata-logo-{v}.png'))

_generate_cdata_logos()

LOGO_DARK  = os.path.join(HERE, 'cdata-logo-dark.png')
LOGO_LIGHT = os.path.join(HERE, 'cdata-logo-light.png')
LOGO_YEL   = os.path.join(HERE, 'cdata-logo-yellow.png')

def _logo_emu_dims(path, target_h=380000):
    from PIL import Image as _I
    pw, ph = _I.open(path).size
    h = Emu(target_h)
    w = int(h * pw / ph)
    return w, h

_CDATA_LOGO_W, _CDATA_LOGO_H = _logo_emu_dims(LOGO_DARK)
_CDATA_LOGO_L = W - ML - _CDATA_LOGO_W
_CDATA_LOGO_T = Emu(450000) - _CDATA_LOGO_H // 2

def add_cdata_logo(slide, variant='dark'):
    path = {'dark': LOGO_DARK, 'light': LOGO_LIGHT, 'yellow': LOGO_YEL}.get(variant, LOGO_DARK)
    slide.shapes.add_picture(path, _CDATA_LOGO_L, _CDATA_LOGO_T, _CDATA_LOGO_W, _CDATA_LOGO_H)

def _generate_qr_code(url):
    import qrcode as _qr
    from PIL import Image as _I
    qr = _qr.QRCode(version=2, box_size=12, border=2,
                     error_correction=_qr.constants.ERROR_CORRECT_M)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color='black', back_color='white').convert('RGBA')
    path = os.path.join(HERE, 'repo-qr.png')
    img.save(path)
    return path

_REPO_QR = _generate_qr_code('https://github.com/JoeKarlsson/connect-ai-agent-demo')


# ── Talk config ───────────────────────────────────────────────────────────────

TALK = {
    'title':    'CData Connect AI — Developer GTM',
    'subtitle': 'A 30-minute look at the audience, the audit, and the plan.',
    'speaker':  'Joe Karlsson',
    'company':  'joekarlsson.com',
    'event':    'CData Final Round',
    'date':     'MAY 2026',
}


# ── CData overrides for pptx_utils layout helpers ────────────────────────────

def overline(slide, text, left=None):
    box(slide, text.upper(), left or ML, Emu(340000), TXT_W, Emu(220000),
        size=11, bold=False, color=CDATA_DARK, font='DM Sans')

def three_cards(slide, items, top=None, card_h=None, animate=True):
    _top    = top    or CONT_S
    _card_h = card_h or Emu(3080000)
    gap = Emu(185000)
    cw  = (CW - 2 * gap) / 3
    pad = Emu(185000)
    groups = []
    for i, item in enumerate(items):
        l = ML + i * (cw + gap)
        bg  = shape(slide, l, _top, cw, _card_h, fill=CDATA_CARD)
        ttl = box(slide, item['title'], l + pad, _top + pad, cw - 2 * pad, Emu(430000),
                  size=17, bold=True, color=CDATA_DARK, wrap=True)
        bt  = _top + pad + Emu(460000)
        bdy = box(slide, item['body'], l + pad, bt, cw - 2 * pad,
                  _card_h - (bt - _top) - pad,
                  size=14, color=CDATA_MID, wrap=True)
        groups.append([bg.shape_id, ttl.shape_id, bdy.shape_id])
    if animate:
        add_appear_animations(slide, groups)

def two_cards(slide, items, top=None, card_h=None, animate=True):
    _top    = top    or CONTENT_T
    _card_h = card_h or Emu(3080000)
    gap = Emu(220000)
    cw  = (CW - gap) / 2
    pad = Emu(185000)
    groups = []
    for i, item in enumerate(items):
        l = ML + i * (cw + gap)
        bg  = shape(slide, l, _top, cw, _card_h, fill=CDATA_CARD)
        ttl = box(slide, item['title'], l + pad, _top + pad, cw - 2 * pad, Emu(430000),
                  size=17, bold=True, color=CDATA_DARK, wrap=True)
        bt  = _top + pad + Emu(460000)
        bdy = box(slide, item['body'], l + pad, bt, cw - 2 * pad,
                  _card_h - (bt - _top) - pad,
                  size=14, color=CDATA_MID, wrap=True)
        groups.append([bg.shape_id, ttl.shape_id, bdy.shape_id])
    if animate:
        add_appear_animations(slide, groups)

def highlight_box(slide, text, top, width=None, height=None):
    import textwrap as _tw
    l = ML
    w = width or CW
    if height is None:
        cpl     = max(20, int((w - Emu(260000)) / 12700 / 7))
        n_lines = max(1, len(_tw.wrap(text, cpl)))
        height  = int(n_lines * int(13 * 12700 * 1.5) + Emu(220000))
    bg   = shape(slide, l, top, w, height, fill=CDATA_CARD, rounded=False)
    bar  = shape(slide, l, top, Emu(60000), height, fill=CDATA_YELLOW, rounded=False)
    body = box(slide, text, l + Emu(170000), top + Emu(100000),
               w - Emu(260000), height - Emu(170000),
               size=13, color=CDATA_DARK, wrap=True)
    return [bg.shape_id, bar.shape_id, body.shape_id]


# ── Speaker scripts ───────────────────────────────────────────────────────────

SCRIPTS = {

    1: """Title slide.
· Introduce yourself briefly — name, the role, one sentence on why you built something before the interview
· Don't explain what you're about to show — just say "before strategy, let me show you something I built this week"
· Transition directly to the demo""",

    2: """Before the strategy.
· Pause here for a beat
· Let the statement land before switching to the demo
· This sets the tone: practitioner, not pitching""",

    3: """Live demo.
· Switch to terminal or play cdata demo.mp4
· Command: python main.py "show me the top 5 customers by total order value"
· Let the agent work visibly — don't rush it
· Watch the tool calls appear: getTables → queryData → format result
· The demo IS the argument — don't narrate it, let it land""",

    4: """Demo result.
· Let this statement breathe
· Don't qualify it — the agent worked, no SQL written, no custom connector
· This is the credibility anchor for everything that follows""",

    5: """Built it. Broke it.
· "I built it" establishes technical credibility
· "Then I broke it" is the pivot to the audit — says you went deeper than a first session
· "Here's what I found" sets up the full presentation
· Short. Let it land. Move on.""",

    6: """Part 1 opens — the developer story.
· Walk through this slowly: the VP asked a question on a Monday. The engineer spent three weeks building the answer.
· "They already had the AI part figured out" — the hard part was never the LLM. It was the data.
· This developer is real. Every engineer who's touched enterprise integrations has lived some version of this.
· The punchline is the last line: Connect AI turns three weeks of integration work into one query.""",

    7: """Cross-source pain.
· This is what the developer is actually trying to do — name it specifically
· "Two different auth models, two different APIs, a maintenance burden they own forever"
· The pain is the maintenance burden, not the initial integration
· "Unless someone has already solved it" — that's the Connect AI setup""",

    8: """Enterprise data slide.
· These four names represent the problem space
· Every senior engineer who's touched enterprise systems knows exactly what these integrations cost
· Don't explain the slide — let the names sit there
· The point: CData has already solved connectivity to all of these""",

    9: """CData's position.
· "Nearly two decades" — not just 350+ sources, but depth, reliability, production hardening
· This is the moat: OAuth edge cases, schema drift, rate limit handling solved across hundreds of API versions
· Nobody replicates this quickly — not Composio, not a team of engineers in 6 months""",

    10: """Not reached yet.
· This is the strategic gap — CData has the product, the audience doesn't know it
· This framing is generous to CData: not "CData failed," but "this audience hasn't been built yet"
· That's the opportunity. That's the job.""",

    11: """Three audiences intro.
· Section divider — pause here, let them reset
· The key insight: these three audiences discover CData in completely different ways
· If you optimize for one, you miss the others
· This is the structural argument for why PMM needs a distinct strategy for each""",

    12: """LLM agents audience.
· LLMs recommend tools when developers ask questions — that recommendation is a distribution channel
· If CData isn't in the training data, isn't in the indexed docs, it doesn't get recommended
· llms.txt is the specific lever — it's what tells AI retrieval systems what Connect AI does
· This is fixable. One file, one week of work.""",

    13: """AI developers audience.
· The builder — the person who is actually writing the code
· They clone repos, hit APIs, and decide in 20 minutes
· GitHub examples, a clear entry point, and comparison content are what move this person
· The quickstart is the most important marketing asset for this audience""",

    14: """Enterprise buyers audience.
· They come in after a developer has already evaluated and used the product
· Sales doesn't get to this buyer until PMM has done its job with the developer
· ROI content, security docs, case studies — these are Sales enablement, not awareness content""",

    15: """Each needs a different approach.
· Section transition slide — bridge to Part 2
· "The audit shows where each one is falling short right now" is the setup
· Don't linger — move directly into the audit""",

    16: """What's actually strong.
· Start with what's working — this is not a gotcha presentation
· First-session time-to-first-API-call is competitive: ~15 minutes
· Enterprise connector depth is real: SAP, NetSuite, Workday, Oracle — Composio doesn't touch these
· MCP tool schemas are well-structured — agents can work with them
· Japanese docs exist — almost no competitor has this
· These are genuine advantages. The signal around them doesn't match.""",

    17: """Complete issue table.
· This is the full scope — 47 issues across 6 categories
· Let them see it before you tell them you're not going to walk through all of it
· "This is what I'm leaving with you" — the full doc is in the repo
· The point: you went deep. 8 hours. Stopwatch running.""",

    18: """Not walking through all 47.
· "I want you to see the scope"
· "Let me show you the ones that matter most strategically"
· These strategic findings are the ones that would surprise leadership — not just annoy a developer""",

    19: """AI documentation standards — what developers expect.
· These five things are table stakes for any AI feature. Every developer building on an AI platform will eventually ask all five.
· Walk the table row by row — don't rush. Each row is both a developer expectation and a specific gap in Connect AI's current docs.
· This slide earns the credibility to say the next three findings aren't just opinion — they're against a known standard.
· Frame it as: "I audited Connect AI against the standard developers actually use when evaluating AI-powered tools." """,

    20: """Finding 1 — invisible.
· "I did what any developer would do — I asked the tool they use"
· This is not a theoretical finding — it's a documented baseline
· The LLM query was run May 8, 2026 — show the data""",

    21: """Screenshot — Composio recommended.
· Let the screenshot land
· Don't editorialize — just let them see it
· "CData not mentioned" is the point""",

    22: """5-query table.
· These are the five queries every developer in this space asks
· Composio appears in all five. CData appears in none.
· This is the discoverability gap in concrete form
· It's not a product problem — it's a training data and content problem""",

    23: """LLMs are where the developer starts.
· This is the mechanism: LLMs train on what's publicly available
· GitHub READMEs, Stack Overflow answers, blog posts with working code — these feed the signal
· "It doesn't exist to them" — not hyperbole. If it doesn't show up in the first place they look, it doesn't exist.""",

    24: """Finding 2 — no agent-specific path.
· Composio picked their target developer and built for them specifically
· "You pick your tool, you get a path built for it" — that's what agent-first onboarding looks like
· This is the content gap, not a product gap""",

    25: """Same seven steps for everyone.
· The Connect AI quickstart exists and it's real — but it's written for the broadest possible audience
· An AI agent builder has to read between the lines to figure out which parts are for them
· That friction is invisible from inside the company but highly visible to the person doing it""",

    26: """CTAs, not docs.
· cdata.com/ai/ is a marketing page — the primary action is signing up, not learning
· A developer landing there to evaluate the product has to go looking for docs
· "Doesn't mean there are no docs" — the docs are real. The entry point isn't agent-first.""",

    27: """Finding 3 — Connect AI effectively orphaned.
· The docs exist. They're complete. They're not the problem.
· The problem is the path to them — and the signal around them
· "Effectively orphaned" is the accurate framing: the product exists, the marketing presence doesn't reflect it""",

    28: """No path to docs.
· Walk through this literally: a developer on cdata.com has to know where to look
· The primary nav doesn't surface Connect AI documentation
· The Connect AI page is a marketing page, not a documentation hub""",

    29: """Wrong place.
· In-product Help links route to CData Sync docs, not Connect AI
· A developer who clicks Help while using Connect AI gets docs for a different product
· This is a specific, fixable bug — but it's also a signal about Connect AI's first-class status""",

    30: """The pattern.
· Three findings, same underlying issue
· Connect AI is hard to find, the experience for AI developers is incomplete, and it's not answering the questions this audience actually has
· This is not a product problem — it's a positioning, content, and distribution problem""",

    31: """Two buckets.
· Not everything from the audit is a PMM problem
· Some of these findings need product conversations — and the best PMM/product work happens earlier
· The point: PMM owns the content and distribution fixes. Product conversations happen with data, not complaints.""",

    32: """PMM fixes vs product conversations.
· PMM fixes ship without waiting: llms.txt rewrite, nav updates, agent-specific quickstart, comparison content, Help link routing
· Product conversations: no sandbox environment, PAT scoping, getInstructions loop, no SQL preview
· "I bring findings, not bugs" — the DX audit is how PMM earns that seat at the product table""",

    33: """90 days.
· Section divider — the plan starts here
· "The audit I just showed you is the document I'd share with product and leadership. Everything in the 90-day plan follows from it."
· This isn't theory — it's execution built on real findings""",

    34: """Days 1-30 framing.
· "Three workstreams. All start on Day 1. Nothing waits for a report."
· Fixes ship the moment they're found
· Data collection and remediation run simultaneously""",

    35: """Three workstreams.
· Funnel instrumentation: map Discover → Evaluate → First API call → Build → Scale. Verify analytics coverage. Pull product usage data.
· Keyword + site health: technical SEO audit, cluster mapping for "enterprise MCP server," "AI agent data access," "SQL connectivity for LLMs"
· LLM discoverability + DX baseline: already have it — show the table. Rewrite llms.txt. Start shipping into training-data surfaces.
· Day 30 deliverable: funnel baseline + keyword cluster map + LLM/DX findings. Data with clear implications, not a strategy deck.""",

    36: """Quick fixes that don't wait.
· These three things don't need a data review — the finding is already documented
· Rewrite the llms.txt preamble: Connect AI product description, quickstart URL, enterprise source list
· Add a docs link to the Connect AI nav
· Fix the in-product Help link routing to CData Sync
· These ship in week one""",

    37: """AI developer entry point.
· "Building with Claude? Start here."
· This is content work, not a product change — one new path through the existing quickstart
· Links from the Connect AI page and the top of the quickstart
· Tested end-to-end with a stopwatch before it ships""",

    38: """Days 31–60.
· Two tracks: one ships Day 31 regardless; the other is driven by the Day 30 funnel data.
· Ships Day 31: re-run the LLM baseline (day 60 vs day 30), publish the Composio vs CData honest comparison piece, fix all 5 AI doc standard misses from the audit.
· Signal-driven fork: discovery bottleneck → SEO + LLM presence content; evaluation drop-off → comparison content + quickstart improvements; activation friction → DX fixes + better error messages.
· Product track: take the 47-issue audit plus baseline data to the product team. Before/after framing — not a wish list.""",

    39: """Days 61–90.
· These are the things we earned the right to build because we spent 60 days gathering real signal.
· Double down on what moved: if the comparison piece moved the needle, add SAP, NetSuite, Workday quickstarts — enterprise sources Composio can't touch. If community seeding moved, conference talk plus OSS working demo.
· Sales enablement: the field can't close technical deals without technical support. Battlecard — Connect AI vs Composio, one honest answer per objection. Plus a "when an engineer asks about Connect AI" one-pager for the field.
· Day 90 re-run: same 5 LLM queries from Day 1. Day 1 result: 0 of 5. Day 90 target: 2–3 of 5. That's the number we're accountable to — not traffic, not MQLs. The metric that proves positioning actually moved.""",

    40: """Goals framework.
· Five metrics. All have a day-30 baseline and a day-90 target.
· LLM discoverability is the leading indicator — if developers are asking LLMs and we're not showing up, nothing else matters
· Time-to-first-API-call is the product funnel metric — under 15 minutes, clean, is the bar
· Community-influenced pipeline is the proof that developer content moves enterprise deals""",

    41: """PMM / DevRel split.
· DevRel earns trust: Discord, GitHub, conference talks, working demos
· PMM turns that trust into pipeline: positioning, comparison content, attribution, Sales enablement
· These are different jobs. They fail when blurred.
· Interface: weekly sync, shared funnel definition. DevRel brings what developers are asking. PMM packages that into content and positioning.""",

    42: """What developer content actually works.
· Specific, runnable, honest about failure modes
· "Seamless enterprise connectivity" without a code sample gets ignored
· "Here's what happens when your PAT expires mid-query and how to handle it" gets bookmarked
· This is the standard for every piece of content in the 90-day plan""",

    43: """What I'm leaving behind.
· Working demo — GitHub repo, public, clone it and run it
· 47-issue DX audit with root causes and recommended fixes
· Composio competitive brief
· LLM baseline data — the before numbers
· Point to the QR code""",

    44: """Close statement.
· Let this land. Don't rush.
· The developer is real. The gap is documented. The fix is clear.
· "That's not a product problem — CData has the product."
· "Getting there is a messaging, content, and distribution problem. That's exactly what I do." """,

    45: """Questions.
· Have dx-audit.md open for specifics
· Have the LLM baseline table ready
· Q&A prep in docs/presentation-structure.md""",
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _section_dark(prs, overline_text, headline, subtext=None):
    """Dark bg section divider — overline top, headline centered."""
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')
    box(s, overline_text.upper(), ML, Emu(340000), TXT_W, Emu(220000),
        size=11, bold=False, color=CDATA_YELLOW, font='DM Sans')
    box(s, headline, ML, Emu(2000000), CW, Emu(2600000),
        size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)
    if subtext:
        box(s, subtext, ML, Emu(4800000), CW, Emu(800000),
            size=20, color=CDATA_GRAY, wrap=True)
    return s

def _statement_dark(prs, text, size=38):
    """Full-screen statement on dark background — text centered."""
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')
    box(s, text, ML, Emu(2400000), CW, Emu(3000000),
        size=size, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)
    return s

def _statement_yellow(prs, text, size=38, subtext=None):
    """Full-screen statement on yellow background — text centered."""
    s = blank(prs)
    set_bg(s, CDATA_YELLOW)
    add_cdata_logo(s, 'yellow')
    box(s, text, ML, Emu(2200000), CW, Emu(3000000),
        size=size, bold=True, color=CDATA_DARK, wrap=True)
    if subtext:
        box(s, subtext, ML, Emu(5200000), CW, Emu(1000000),
            size=18, color=CDATA_DARK, wrap=True)
    return s

def _statement_light(prs, text, size=38, subtext=None, overline_text=None):
    """Full-screen statement on light background — text centered."""
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    if overline_text:
        overline(s, overline_text)
    box(s, text, ML, Emu(2200000), CW, Emu(3000000),
        size=size, bold=True, color=CDATA_DARK, wrap=True)
    if subtext:
        box(s, subtext, ML, Emu(5200000), CW, Emu(1000000),
            size=18, color=CDATA_MID, wrap=True)
    return s


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = blank(prs)
    set_bg(s, RGBColor(0xFF, 0xFF, 0xFF))
    shape(s, LP_L, LP_T, LP_W, LP_H, fill=CDATA_YELLOW, rounded=False)
    add_cdata_logo(s, 'dark')
    box(s, f"{TALK['event'].upper()}  ·  {TALK['date']}", ML, Emu(1280000), TXT_W, Emu(300000),
        size=12, bold=True, color=CDATA_DARK, font='DM Sans')
    box(s, TALK['title'], ML, Emu(1680000), TXT_W, Emu(1600000),
        size=44, bold=True, color=CDATA_DARK, wrap=True)
    box(s, TALK['subtitle'], ML, Emu(3450000), TXT_W, Emu(900000),
        size=17, color=CDATA_GRAY, wrap=True)
    box(s, TALK['speaker'], ML, Emu(4700000), TXT_W, Emu(370000),
        size=16, bold=True, color=CDATA_DARK)
    box(s, TALK['company'], ML, Emu(5100000), TXT_W, Emu(330000),
        size=13, color=CDATA_GRAY, url='https://joekarlsson.com')
    set_notes(s, SCRIPTS[1])


def slide_02_before_strategy(prs):
    s = _statement_dark(prs,
        'Before the strategy — let me show you what I built.',
        size=44)
    set_notes(s, SCRIPTS[2])


def slide_03_demo(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    mid_y = H // 2 - Emu(900000)
    box(s, '▶', 0, mid_y, W, Emu(700000),
        size=44, bold=True, color=CDATA_YELLOW, align=PP_ALIGN.CENTER)
    box(s, 'LIVE DEMO', 0, mid_y + Emu(750000), W, Emu(400000),
        size=14, bold=True, color=CDATA_YELLOW, align=PP_ALIGN.CENTER, font='DM Sans')
    box(s, 'Connect AI Agent', 0, mid_y + Emu(1200000), W, Emu(500000),
        size=28, bold=True, color=RGBColor(0xF0, 0xF0, 0xF0), align=PP_ALIGN.CENTER)

    cmd_y = mid_y + Emu(1850000)
    cmd_h = Emu(480000)
    shape(s, ML + Emu(800000), cmd_y, CW - Emu(1600000), cmd_h,
          fill=RGBColor(0x2A, 0x2A, 0x2A), rounded=True)
    box(s, 'python main.py "show me the top 5 customers by total order value"',
        ML + Emu(1000000), cmd_y + Emu(120000), CW - Emu(2000000), Emu(270000),
        size=13, color=CDATA_YELLOW, font='Courier New', wrap=False)
    set_notes(s, SCRIPTS[3])


def slide_04_demo_result(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, 'That agent just queried live enterprise data through Connect AI\'s MCP server.',
        ML, Emu(2000000), CW, Emu(1800000),
        size=38, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)
    box(s, 'No custom connector. No SQL written by hand.',
        ML, Emu(4000000), CW, Emu(700000),
        size=28, color=CDATA_YELLOW, wrap=True)
    set_notes(s, SCRIPTS[4])


def slide_05_built_and_broke(prs):
    s = blank(prs)
    set_bg(s, CDATA_YELLOW)
    add_cdata_logo(s, 'yellow')

    box(s, 'I built it.', ML, Emu(1400000), CW, Emu(900000),
        size=56, bold=True, color=CDATA_DARK, wrap=True)
    box(s, 'Then I broke it.', ML, Emu(2400000), CW, Emu(900000),
        size=56, bold=True, color=CDATA_DARK, wrap=True)
    box(s, "Here's what I found.", ML, Emu(3500000), CW, Emu(700000),
        size=32, color=CDATA_DARK, wrap=True)
    set_notes(s, SCRIPTS[5])


# ── PART 1: The AI Developer Opportunity ─────────────────────────────────────

def slide_06_developer_who_needs(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Part 1: The Opportunity')

    # Main story — centered vertically
    box(s, 'There is a developer right now building an AI agent that needs to query Salesforce and NetSuite together.',
        ML, Emu(1400000), CW, Emu(1400000),
        size=30, bold=True, color=CDATA_DARK, wrap=True)

    box(s, 'They Googled it. They asked Claude. They checked GitHub.',
        ML, Emu(3000000), CW, Emu(700000),
        size=24, color=CDATA_MID, wrap=True)

    box(s, 'Composio came up. Custom connector tutorials came up.',
        ML, Emu(3800000), CW, Emu(600000),
        size=24, color=CDATA_MID, wrap=True)

    box(s, 'CData didn\'t come up.',
        ML, Emu(4550000), CW, Emu(580000),
        size=28, bold=True, color=CDATA_DARK, wrap=True)

    shape(s, ML, Emu(5280000), CW, Emu(34000), fill=CDATA_YELLOW, rounded=False)
    box(s, 'CData Connect AI is exactly what they need. They just can\'t find it.',
        ML, Emu(5380000), CW, Emu(480000),
        size=16, color=CDATA_DARK, wrap=True)

    set_notes(s, SCRIPTS[6])


def slide_07_cross_source_pain(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'They want their AI agent to answer questions about Salesforce pipeline and NetSuite financials in one query.',
        ML, Emu(1600000), CW, Emu(1800000),
        size=30, bold=True, color=CDATA_DARK, wrap=True)

    box(s, 'Two different auth models.  Two different APIs.  A maintenance burden they own forever.',
        ML, Emu(3600000), CW, Emu(800000),
        size=22, color=CDATA_MID, wrap=True)

    box(s, 'Unless someone has already solved it.',
        ML, Emu(4500000), CW, Emu(600000),
        size=22, bold=True, color=CDATA_DARK, wrap=True)
    set_notes(s, SCRIPTS[7])


def slide_08_enterprise_data(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    for i, name in enumerate(['Salesforce.', 'SAP.', 'Workday.', 'NetSuite.']):
        box(s, name, ML, Emu(1000000) + i * Emu(1200000), CW, Emu(1000000),
            size=52, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=False)

    box(s, 'The data that runs their business, connected to the AI agents they\'re building.',
        ML, Emu(5800000), CW, Emu(500000),
        size=15, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[8])


def slide_09_cdata_position(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'CData has nearly two decades of enterprise connector experience.',
        ML, Emu(1900000), CW, Emu(1600000),
        size=40, bold=True, color=CDATA_DARK, wrap=True)
    box(s, 'More sources, more security, more production reliability than anyone else in this space.',
        ML, Emu(3700000), CW, Emu(900000),
        size=22, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[9])


def slide_10_not_reached(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'CData built that reputation talking to enterprise buyers and the driver community.',
        ML, Emu(1700000), CW, Emu(1400000),
        size=32, bold=True, color=CDATA_DARK, wrap=True)

    shape(s, ML, Emu(3400000), Emu(120000), Emu(120000), fill=CDATA_YELLOW, rounded=False)
    box(s, 'This AI developer audience hasn\'t been reached yet.',
        ML + Emu(300000), Emu(3370000), CW - Emu(300000), Emu(700000),
        size=28, bold=True, color=CDATA_DARK, wrap=True)

    box(s, 'That\'s the opportunity.',
        ML, Emu(4250000), CW, Emu(600000),
        size=22, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[10])


def slide_11_three_audiences(prs):
    s = _section_dark(prs, 'The Audience',
        'Three distinct audiences. Three completely different ways they find us.',
        'Getting this right is the structural problem PMM needs to solve.')
    set_notes(s, SCRIPTS[11])


def slide_12_llm_agents(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Audience 1')
    box(s, 'LLM Agents', ML, Emu(620000), CW, Emu(900000),
        size=52, bold=True, color=CDATA_DARK, wrap=False)
    box(s, 'Claude. Cursor. Copilot.', ML, Emu(1620000), CW, Emu(500000),
        size=20, color=CDATA_GRAY, wrap=True)

    items = [
        'Structured data they can read and index',
        'A useful llms.txt — not a sitemap, a product guide',
        'Agent-readable docs so they recommend CData when a developer asks for help',
    ]
    dy = Emu(2600000)
    for item in items:
        shape(s, ML, dy + Emu(140000), Emu(90000), Emu(90000), fill=CDATA_YELLOW, rounded=False)
        box(s, item, ML + Emu(260000), dy, CW - Emu(260000), Emu(420000),
            size=20, color=CDATA_DARK, wrap=True)
        dy += Emu(520000)

    box(s, 'They discover CData through training data, indexed docs, and structured content.',
        ML, Emu(4500000), CW, Emu(400000),
        size=14, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[12])


def slide_13_ai_developers(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Audience 2')
    box(s, 'AI Developers', ML, Emu(620000), CW, Emu(900000),
        size=52, bold=True, color=CDATA_DARK, wrap=False)
    box(s, 'The builders. They clone repos, hit APIs, and decide in 20 minutes.',
        ML, Emu(1620000), CW, Emu(500000),
        size=20, color=CDATA_GRAY, wrap=True)

    items = [
        'GitHub examples they can clone and run',
        'A clear entry point built for agent builders, not platform evaluators',
        'Comparison content that earns the recommendation',
    ]
    dy = Emu(2600000)
    for item in items:
        shape(s, ML, dy + Emu(140000), Emu(90000), Emu(90000), fill=CDATA_YELLOW, rounded=False)
        box(s, item, ML + Emu(260000), dy, CW - Emu(260000), Emu(420000),
            size=20, color=CDATA_DARK, wrap=True)
        dy += Emu(520000)

    box(s, 'They discover CData through search, LLM recommendations, and GitHub.',
        ML, Emu(4500000), CW, Emu(400000),
        size=14, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[13])


def slide_14_enterprise_buyers(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Audience 3')
    box(s, 'Enterprise Buyers', ML, Emu(620000), CW, Emu(900000),
        size=52, bold=True, color=CDATA_DARK, wrap=False)
    box(s, 'They come in after a developer evaluation. Sales reaches them through PMM.',
        ML, Emu(1620000), CW, Emu(500000),
        size=20, color=CDATA_GRAY, wrap=True)

    items = [
        'ROI content and security / compliance documentation',
        'Case studies with numbers a VP can present to a committee',
        'Sales enablement that closes — built from what the developer already evaluated',
    ]
    dy = Emu(2600000)
    for item in items:
        shape(s, ML, dy + Emu(140000), Emu(90000), Emu(90000), fill=CDATA_YELLOW, rounded=False)
        box(s, item, ML + Emu(260000), dy, CW - Emu(260000), Emu(420000),
            size=20, color=CDATA_DARK, wrap=True)
        dy += Emu(520000)

    box(s, 'They discover CData through the developers already using it.',
        ML, Emu(4500000), CW, Emu(400000),
        size=14, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[14])


def slide_15_different_approach(prs):
    s = _section_dark(prs, 'The Bridge',
        'Each one needs a different approach.',
        'The audit shows where each one is falling short right now.')
    set_notes(s, SCRIPTS[15])


# ── PART 2: The Audit ─────────────────────────────────────────────────────────

def slide_16_whats_working(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Part 2: The Audit')
    box(s, 'Before what needs work — here\'s what\'s actually strong.',
        ML, Emu(620000), CW, Emu(900000),
        size=36, bold=True, color=CDATA_DARK, wrap=True)

    strengths = [
        ('~15 min to first API call', 'Competitive. Composio is faster to hello-world but shallower.'),
        ('Enterprise connector depth', 'SAP, NetSuite, Workday, Oracle. Composio barely touches these.'),
        ('MCP tool schemas', 'Well-structured. Agents can work with them out of the box.'),
        ('Japanese documentation', 'Almost no competitor has this. Real enterprise signal.'),
    ]
    dy = Emu(1900000)
    for title, body in strengths:
        shape(s, ML, dy, CW, Emu(520000), fill=CDATA_CARD, rounded=False)
        shape(s, ML, dy, Emu(10000), Emu(520000), fill=CDATA_YELLOW, rounded=False)
        box(s, title, ML + Emu(220000), dy + Emu(80000), Emu(3600000), Emu(300000),
            size=16, bold=True, color=CDATA_DARK, wrap=False)
        box(s, body, ML + Emu(3900000), dy + Emu(80000), CW - Emu(4100000), Emu(360000),
            size=14, color=CDATA_MID, wrap=True)
        dy += Emu(580000)

    box(s, 'The audit isn\'t saying the product is broken. It\'s saying the signal around it doesn\'t reflect what the product actually is.',
        ML, Emu(4380000), CW, Emu(500000),
        size=13, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[16])


def slide_17_issue_table(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, 'THE AUDIT', ML, Emu(330000), CW, Emu(180000),
        size=11, bold=True, color=CDATA_YELLOW, wrap=False)
    box(s, '47 issues. 6 categories. First session only.',
        ML, Emu(530000), CW, Emu(560000),
        size=34, bold=True, color=CDATA_YELLOW, wrap=True)
    box(s, 'Quickstart (6)  ·  Production Readiness (12)  ·  Agent & MCP (6)  ·  Docs & Navigation (15)  ·  SEO (8)',
        ML, Emu(1120000), CW, Emu(200000),
        size=11, color=CDATA_GRAY, wrap=False)
    shape(s, ML, Emu(1360000), CW, Emu(4000), fill=CDATA_DIVIDER, rounded=False)

    issues = [
        "Wrong Accept header breaks every copy-pasted implementation",
        "Setup Guide shows navigation steps you already completed",
        "401 error doesn't tell you the username is wrong",
        "No terminal-first setup path — dashboard required for OAuth",
        "Tool responses embed verbose schema metadata on every call",
        "No virtual environment guidance in quickstart",
        "Rate limit: shared per-user ceiling, absent from MCP docs",
        "Cross-source JOINs pull full result sets before filtering",
        "Query timeout is 30s with GUI-only override — no API param",
        "PATs have no scoping — one token, all sources, read/write",
        "Log retention is 7 days, dashboard-only, no export API",
        '"Build your own MCP tools" gated to Business tier',
        "SSO and RBAC are Business-tier-only — no lower-tier controls",
        "Local MCP server is stdio-only — no remote deployments",
        "Error codes are numeric 0–79 with no published reference",
        "No sandbox — every query hits production credentials",
        "100M record/month fair use limit is undefined",
        "Per-framework docs are siloed — no unified compat matrix",
        "getInstructions routes agents into expensive metadata loop",
        "Three-part table names required, not surfaced until failure",
        "Tool naming breaks for multi-word or underscore source names",
        "No result streaming or pagination through MCP interface",
        "DataType enum table in REST docs, not linked from MCP pages",
        "MCP config examples store credentials without secrets guidance",
        "Docs fragmented across subdomains — domain authority split",
        "Docs absent from primary navigation — no docs link in header",
        "No model or version disclosure for any AI feature",
        "No latency specifications — fitness for use case unknown",
        "No accuracy or limitation disclosures — AI is described vaguely",
        "No developer override or control mechanism documented",
        "No cost or token implications disclosed for AI features",
        "docs.cloud.cdata.com fails agent-readiness checks",
        "In-product Help links resolve to CData Sync, not Connect AI",
        "Quick Start has no prerequisites section",
        "Zero screenshots in the Quick Start Guide",
        "Code samples are missing imports — not copy-pasteable",
        "No error handling in any code sample",
        "FAQ has 6 questions, no cross-links — nearly useless",
        "Docs home has no beginner/advanced split — navigation is flat",
        "No JSON-LD structured data on any audited page",
        "Open Graph and Twitter Card tags unconfirmed",
        "Docs home has no H1 — page starts with H2",
        "No last-updated dates on docs pages",
        "Blog posts lack BlogPosting schema — no author/date signals",
        "Sitemap is manually generated — risks going stale",
        "hreflang tags absent despite Japanese documentation existing",
        "Image alt text uses filenames — keyword opportunity missed",
    ]

    COL_W  = Emu(3380000)
    COL_GAP = Emu(285000)
    ITEM_H  = Emu(193000)
    NUM_W   = Emu(290000)
    TOP     = Emu(1480000)
    TEXT_COLOR = RGBColor(0xBB, 0xBB, 0xB4)

    splits = [(0, 16, ML), (16, 32, ML + COL_W + COL_GAP), (32, 47, ML + 2*(COL_W + COL_GAP))]
    for start, end, col_x in splits:
        dy = TOP
        for i, issue in enumerate(issues[start:end]):
            num = start + i + 1
            box(s, f'{num:02d}', col_x, dy, NUM_W, ITEM_H,
                size=8, bold=True, color=CDATA_YELLOW, wrap=False)
            box(s, issue, col_x + NUM_W, dy, COL_W - NUM_W, ITEM_H,
                size=8, color=TEXT_COLOR, wrap=False)
            dy += ITEM_H

    source_label(s, 'Full audit: docs/dx-audit.md  ·  github.com/JoeKarlsson/connect-ai-agent-demo')
    set_notes(s, SCRIPTS[17])


def slide_18_not_all_47(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, "I'm not going to walk you through all 47.",
        ML, Emu(2000000), CW, Emu(1200000),
        size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)
    box(s, 'I want you to see the scope.',
        ML, Emu(3400000), CW, Emu(700000),
        size=26, color=CDATA_YELLOW, wrap=True)
    box(s, 'Let me show you the ones that matter most strategically.',
        ML, Emu(4200000), CW, Emu(700000),
        size=22, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[18])


def slide_19_ai_doc_standards(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'The Standard')
    box(s, 'AI developers expect five things from AI feature documentation.',
        ML, Emu(620000), CW, Emu(800000),
        size=32, bold=True, color=CDATA_DARK, wrap=True)
    box(s, 'CData Connect AI provides none of them.',
        ML, Emu(1520000), CW, Emu(500000),
        size=22, bold=True, color=CDATA_DARK, wrap=True)

    rows = [
        ('Model + version',
         'Which model powers the feature? What version? Will you notify developers when it changes?',
         'Not disclosed anywhere in product or docs.'),
        ('Input/output specs',
         'What goes in, what comes out, what are the token limits? If certain inputs produce worse results, say so.',
         'No token limit disclosure. No input constraints documented.'),
        ('Accuracy + limitations',
         '"May produce incorrect results for complex nested queries" is useful. "Results may vary" is not.',
         '"98.5% accuracy across 378 queries" — no methodology, no connector breakdown.'),
        ('Failure modes + recovery',
         'What happens when the AI gets it wrong? Confidence score? Fallback? Undo?',
         'Black box. No fallback, no dry run, no SQL preview before write operations.'),
        ('Cost implications',
         'If the feature consumes tokens or compute, map it to the developer\'s bill. Give them a formula.',
         '/connect-ai/pricing returns a 404. No token implications disclosed anywhere.'),
    ]

    dy = Emu(2200000)
    row_h = Emu(850000)
    for i, (standard, expectation, finding) in enumerate(rows):
        fill = CDATA_CARD if i % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF5)
        shape(s, ML, dy, CW, row_h - Emu(30000), fill=fill, rounded=False)
        shape(s, ML, dy, Emu(10000), row_h - Emu(30000), fill=CDATA_YELLOW, rounded=False)
        box(s, standard, ML + Emu(200000), dy + Emu(100000),
            Emu(2600000), Emu(300000),
            size=13, bold=True, color=CDATA_DARK, wrap=False)
        box(s, finding, ML + Emu(2900000), dy + Emu(100000),
            CW - Emu(3100000), row_h - Emu(200000),
            size=12, color=RGBColor(0xC0, 0x30, 0x30), wrap=True)
        dy += row_h

    set_notes(s, SCRIPTS[19])


def slide_20_finding1_intro(prs):  # noqa: F811
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Strategic Finding #1')
    box(s, 'CData Is Invisible in the AI Developer\'s Information Environment.',
        ML, Emu(620000), CW, Emu(1200000),
        size=36, bold=True, color=CDATA_DARK, wrap=True)
    box(s, 'I asked Claude Code which tool to use to connect AI agents to enterprise data.',
        ML, Emu(2100000), CW, Emu(800000),
        size=24, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[20])


def slide_21_screenshot(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    # Simulated LLM response card
    card_t = Emu(1200000)
    card_h = Emu(4200000)
    shape(s, ML, card_t, CW, card_h, fill=RGBColor(0x28, 0x28, 0x28), rounded=True)
    shape(s, ML, card_t, CW, Emu(80000), fill=RGBColor(0x40, 0x40, 0x40), rounded=False)

    box(s, 'Claude Code — May 8, 2026',
        ML + Emu(300000), card_t + Emu(150000), CW, Emu(300000),
        size=11, color=CDATA_GRAY, wrap=False)

    response = (
        '"For connecting AI agents to enterprise data via MCP, I recommend Composio. '
        'It supports 250+ integrations including Salesforce, Jira, GitHub, Slack, and HubSpot '
        'with a clean MCP server implementation. '
        'The free tier covers most development needs and the documentation is excellent for agent builders.\n\n'
        'For enterprise data specifically, you might also look at Airbyte for batch pipelines '
        'or build custom MCP servers for your specific sources."'
    )
    box(s, response, ML + Emu(300000), card_t + Emu(600000), CW - Emu(600000), Emu(3400000),
        size=15, color=RGBColor(0xE0, 0xE0, 0xE0), wrap=True, font='Courier New')

    box(s, 'CData not mentioned.',
        ML, card_t + card_h + Emu(300000), CW, Emu(500000),
        size=24, bold=True, color=CDATA_YELLOW, wrap=False)
    set_notes(s, SCRIPTS[21])


def slide_21_5_queries(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'LLM Baseline — May 8, 2026')
    box(s, 'I ran 5 key queries. CData appeared in 0 of them.',
        ML, Emu(620000), CW, Emu(800000),
        size=34, bold=True, color=CDATA_DARK, wrap=True)

    queries = [
        ('Connect AI agents to enterprise data',    'Composio (recommended)',        '—'),
        ('MCP server for enterprise data sources',  'Composio, custom build',        '—'),
        ('SQL interface for AI agents',             'LangChain, Composio',           '—'),
        ('AI agent data connectivity platform',     'Composio — "best overall"',     '—'),
        ('Connect Salesforce + NetSuite to LLM',    'Composio, custom API clients',  '—'),
    ]

    header_y = Emu(1700000)
    row_h    = Emu(760000)

    # Header
    shape(s, ML, header_y, CW, Emu(340000), fill=CDATA_DARK, rounded=False)
    cols = [(0, 'Query'), (Emu(4600000), 'Composio'), (Emu(8600000), 'CData')]
    for off, hdr in cols:
        box(s, hdr, ML + off + Emu(180000), header_y + Emu(70000), Emu(3800000), Emu(240000),
            size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=False)

    for i, (q, composio, cdata) in enumerate(queries):
        row_y = header_y + Emu(340000) + i * row_h
        fill  = CDATA_CARD if i % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF5)
        shape(s, ML, row_y, CW, row_h - Emu(40000), fill=fill, rounded=False)
        box(s, q,       ML + Emu(180000),              row_y + Emu(170000), Emu(4200000), Emu(420000),
            size=13, color=CDATA_DARK, wrap=True)
        box(s, composio, ML + Emu(4600000) + Emu(180000), row_y + Emu(170000), Emu(3800000), Emu(420000),
            size=13, color=CDATA_DARK, wrap=True)
        box(s, cdata,   ML + Emu(8600000) + Emu(180000), row_y + Emu(170000), Emu(1600000), Emu(420000),
            size=13, color=CDATA_GRAY, wrap=False)

    set_notes(s, SCRIPTS[22])


def slide_22_llms_where_dev_starts(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, 'LLMs are where this developer starts.',
        ML, Emu(1700000), CW, Emu(1100000),
        size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)
    box(s, 'If CData isn\'t in the training data, isn\'t in the GitHub repos they clone, isn\'t in the Stack Overflow answers they find —',
        ML, Emu(3100000), CW, Emu(1000000),
        size=22, color=CDATA_GRAY, wrap=True)
    box(s, "it doesn't exist to them.",
        ML, Emu(4250000), CW, Emu(700000),
        size=28, bold=True, color=CDATA_YELLOW, wrap=True)
    set_notes(s, SCRIPTS[23])


def slide_23_finding2_no_agent_path(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Strategic Finding #2')
    box(s, 'No Agent-Specific Path.',
        ML, Emu(1400000), CW, Emu(900000),
        size=44, bold=True, color=CDATA_DARK, wrap=True)

    box(s, 'Composio has dedicated quickstarts for Claude, ChatGPT, Cursor, and Copilot.',
        ML, Emu(2600000), CW, Emu(700000),
        size=24, color=CDATA_DARK, wrap=True)
    box(s, 'You pick your tool. You get a path built for it.',
        ML, Emu(3450000), CW, Emu(600000),
        size=24, bold=True, color=CDATA_DARK, wrap=True)
    set_notes(s, SCRIPTS[24])


def slide_24_same_seven_steps(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'The Connect AI quickstart covers AI agent builders, BI analysts, and platform evaluators in the same seven steps.',
        ML, Emu(1800000), CW, Emu(1600000),
        size=32, bold=True, color=CDATA_DARK, wrap=True)
    box(s, "It's broad by design. But a developer building an agent has to figure out which parts are for them.",
        ML, Emu(3700000), CW, Emu(900000),
        size=22, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[25])


def slide_25_cta_not_docs(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'cdata.com/ai/ is a marketing page.',
        ML, Emu(1900000), CW, Emu(900000),
        size=38, bold=True, color=CDATA_DARK, wrap=True)
    box(s, 'The primary CTAs go to a signup form.',
        ML, Emu(3000000), CW, Emu(600000),
        size=24, color=CDATA_MID, wrap=True)
    box(s, "The docs and quickstart aren't above the fold — a developer landing there to evaluate the product has to go looking.",
        ML, Emu(3800000), CW, Emu(900000),
        size=20, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[26])


def slide_26_finding3_orphaned(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Strategic Finding #3')
    box(s, 'Connect AI Is Not a First-Class Product on the Marketing Site.',
        ML, Emu(620000), CW, Emu(1000000),
        size=36, bold=True, color=CDATA_DARK, wrap=True)

    box(s, 'Connect AI has its own documentation.',
        ML, Emu(2000000), CW, Emu(500000),
        size=24, color=CDATA_DARK, wrap=True)
    box(s, "It's real, it's complete, and it's effectively orphaned.",
        ML, Emu(2650000), CW, Emu(600000),
        size=24, bold=True, color=CDATA_DARK, wrap=True)
    set_notes(s, SCRIPTS[27])


def slide_27_no_path_to_docs(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'A developer landing on cdata.com trying to find Connect AI documentation has to know where to look.',
        ML, Emu(1800000), CW, Emu(1600000),
        size=32, bold=True, color=CDATA_DARK, wrap=True)
    box(s, "It's not in the primary nav. The Connect AI page is a marketing page with a signup CTA.",
        ML, Emu(3600000), CW, Emu(900000),
        size=22, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[28])


def slide_28_wrong_place(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, 'When an AI developer clicks Help inside Connect AI, they end up in CData Sync documentation.',
        ML, Emu(1000000), CW, Emu(1600000),
        size=30, bold=True, color=CDATA_DARK, wrap=True)

    shape(s, ML, Emu(2900000), Emu(2800000), Emu(560000), fill=CDATA_CARD, rounded=False)
    shape(s, ML, Emu(2900000), Emu(10000), Emu(560000), fill=CDATA_YELLOW, rounded=False)
    box(s, 'The docs exist.', ML + Emu(220000), Emu(2980000), Emu(2600000), Emu(400000),
        size=20, bold=True, color=CDATA_DARK, wrap=False)

    shape(s, ML + Emu(3000000), Emu(2900000), Emu(2800000), Emu(560000), fill=CDATA_CARD, rounded=False)
    shape(s, ML + Emu(3000000), Emu(2900000), Emu(10000), Emu(560000), fill=CDATA_DARK, rounded=False)
    box(s, "The path to them doesn't.", ML + Emu(3220000), Emu(2980000), Emu(2600000), Emu(400000),
        size=20, bold=True, color=CDATA_DARK, wrap=False)
    set_notes(s, SCRIPTS[29])


def slide_29_the_pattern(prs):
    s = _section_dark(prs, 'The Pattern',
        'Connect AI is hard to find, the experience for AI developers is incomplete, and it\'s not answering the questions this audience actually has.',
        'Three findings. Same underlying issue.')
    set_notes(s, SCRIPTS[30])


def slide_30_two_buckets(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'What Comes Next')
    box(s, 'These findings split into two buckets.',
        ML, Emu(620000), CW, Emu(900000),
        size=40, bold=True, color=CDATA_DARK, wrap=True)

    two_cards(s, [
        {
            'title': 'PMM fixes',
            'body':  (
                'llms.txt rewrite\n'
                'Add docs link to Connect AI nav\n'
                'Agent-specific quickstart path\n'
                'Comparison content (vs Composio)\n'
                'Fix in-product Help link routing'
            ),
        },
        {
            'title': 'Product conversations',
            'body':  (
                'No sandbox environment\n'
                'PAT scoping (no read-only option)\n'
                'getInstructions metadata loop\n'
                'No SQL preview before write operations\n'
                'stdio-only local MCP server'
            ),
        },
    ], top=Emu(1900000), card_h=Emu(3600000), animate=False)
    set_notes(s, SCRIPTS[31])


def slide_31_pmm_product_bridge(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')

    box(s, "The point isn't \"I bring bugs to product.\"",
        ML, Emu(1700000), CW, Emu(900000),
        size=36, bold=True, color=CDATA_DARK, wrap=True)
    box(s, 'The best PMM/product work happens earlier — helping frame what\'s being built before it ships, technically reviewing it, and making sure the external story matches what developers actually experience.',
        ML, Emu(2900000), CW, Emu(1400000),
        size=20, color=CDATA_MID, wrap=True)
    box(s, 'The DX audit is an example of all three.',
        ML, Emu(4550000), CW, Emu(600000),
        size=22, bold=True, color=CDATA_DARK, wrap=True)
    set_notes(s, SCRIPTS[32])


# ── PART 3: The GTM Plan ──────────────────────────────────────────────────────

def slide_32_90_days(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, 'Part 3: The GTM Plan', 0, Emu(800000), W, Emu(400000),
        size=13, bold=True, color=CDATA_GRAY, align=PP_ALIGN.CENTER, font='DM Sans')
    box(s, '90 days.', 0, Emu(1900000), W, Emu(1800000),
        size=100, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    box(s, 'The audit I just showed you is the document I\'d share with product and leadership. Everything in this plan follows from it.',
        ML, Emu(4200000), CW, Emu(700000),
        size=16, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[33])


def slide_33_days_130_intro(prs):
    s = _statement_yellow(prs,
        'Days 1–30: Establish the baseline. Ship the quick fixes. Run everything in parallel.',
        size=34)
    set_notes(s, SCRIPTS[34])


def slide_34_three_workstreams(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Days 1–30')
    box(s, 'Three parallel workstreams. All start on Day 1. Nothing waits for a report.',
        ML, Emu(620000), CW, Emu(900000),
        size=30, bold=True, color=CDATA_DARK, wrap=True)

    three_cards(s, [
        {
            'title': 'Funnel Instrumentation',
            'body':  (
                'Map Discover → Evaluate → First API call → Build → Scale.\n\n'
                'Verify analytics coverage at each stage. Pull product usage data — '
                'which connectors are getting trial activity? Where are developers dropping off?\n\n'
                'First funnel health report by week 2.'
            ),
        },
        {
            'title': 'Keyword + Site Health',
            'body':  (
                'Technical SEO audit: redirect chains, Core Web Vitals, crawl health, internal linking.\n\n'
                'Keyword clusters: "enterprise MCP server," "AI agent data access," '
                '"SQL connectivity for LLMs."\n\n'
                'Remediate obvious issues immediately — don\'t wait for the full audit.'
            ),
        },
        {
            'title': 'LLM Discoverability + DX',
            'body':  (
                'LLM baseline already run: CData absent from all 5 key queries.\n\n'
                'llms.txt rewrite: Connect AI product guide, not a page index. One week.\n\n'
                'Start shipping into training-data surfaces in parallel: README improvements, '
                'Stack Overflow, GitHub examples.'
            ),
        },
    ], top=Emu(1900000), card_h=Emu(4700000), animate=False)

    set_notes(s, SCRIPTS[35])


def slide_35_immediate_fixes(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Days 1–30')
    box(s, "Immediate fixes that don't wait for data.",
        ML, Emu(620000), CW, Emu(800000),
        size=36, bold=True, color=CDATA_DARK, wrap=True)

    fixes = [
        ('Rewrite llms.txt', 'Lead with a Connect AI product description: what it does, how it differs from Composio, the quickstart URL, the enterprise source list. Currently a page index. One week of work.'),
        ('Add a docs link to the Connect AI nav', 'The documentation exists. Put it where the developer can find it. This is a nav item, not a campaign.'),
        ('Fix the in-product Help link', 'Help currently routes to CData Sync documentation. Route it to Connect AI. This is a single redirect.'),
    ]

    dy = Emu(1800000)
    for i, (title, body) in enumerate(fixes):
        shape(s, ML, dy, CW, Emu(1100000), fill=CDATA_CARD, rounded=False)
        shape(s, ML, dy, Emu(10000), Emu(1100000), fill=CDATA_YELLOW, rounded=False)
        box(s, title, ML + Emu(240000), dy + Emu(100000), Emu(4000000), Emu(340000),
            size=17, bold=True, color=CDATA_DARK, wrap=False)
        box(s, body, ML + Emu(240000), dy + Emu(480000), CW - Emu(440000), Emu(540000),
            size=13, color=CDATA_MID, wrap=True)
        dy += Emu(1200000)
    set_notes(s, SCRIPTS[36])


def slide_36_entry_point(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Days 1–30')
    box(s, 'AI developer entry point.',
        ML, Emu(620000), CW, Emu(800000),
        size=44, bold=True, color=CDATA_DARK, wrap=True)

    # Simulated quickstart path
    shape(s, ML, Emu(1900000), CW, Emu(1600000), fill=CDATA_CARD, rounded=True)
    box(s, '"Building with Claude? Start here."',
        ML + Emu(400000), Emu(2200000), CW - Emu(800000), Emu(600000),
        size=26, bold=True, color=CDATA_DARK, wrap=True, font='Courier New')

    box(s, 'Build an agent-specific path into the existing quickstart. Link it from the Connect AI page.',
        ML, Emu(3800000), CW, Emu(700000),
        size=20, color=CDATA_MID, wrap=True)
    box(s, 'This is content work, not a product change. Tested with a stopwatch before it ships.',
        ML, Emu(4600000), CW, Emu(600000),
        size=18, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[37])


def slide_37_days_3160(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Days 31–60')
    box(s, 'Two tracks. One ships Day 31 — the other follows data.',
        ML, Emu(620000), CW, Emu(500000),
        size=32, bold=True, color=CDATA_DARK, wrap=True)

    COL1_W = Emu(4900000)
    COL2_X = ML + Emu(5300000)
    COL2_W = CW - Emu(5300000)
    TOP = Emu(1350000)

    # Left column: ships regardless
    shape(s, ML, TOP, COL1_W, Emu(280000), fill=CDATA_YELLOW, rounded=False)
    box(s, 'Ships Day 31 — no signal needed', ML + Emu(120000), TOP + Emu(50000),
        COL1_W - Emu(240000), Emu(200000), size=13, bold=True, color=CDATA_DARK, wrap=False)

    left_items = [
        'Re-run LLM baseline — day 60 vs day 30',
        'Composio vs CData honest comparison (the "which should I use?" piece)',
        'Fix all 5 AI doc standard misses from the audit',
    ]
    dy = TOP + Emu(360000)
    for item in left_items:
        shape(s, ML, dy + Emu(60000), Emu(8000), Emu(200000), fill=CDATA_DARK, rounded=False)
        box(s, item, ML + Emu(180000), dy, COL1_W - Emu(180000), Emu(320000),
            size=15, color=CDATA_DARK, wrap=True)
        dy += Emu(400000)

    # Vertical divider
    shape(s, ML + Emu(5100000), TOP, Emu(5000), Emu(1600000), fill=CDATA_DIVIDER, rounded=False)

    # Right column: signal-driven fork
    shape(s, COL2_X, TOP, COL2_W, Emu(280000), fill=CDATA_CARD, rounded=False)
    box(s, 'Day 30 data decides', COL2_X + Emu(80000), TOP + Emu(50000),
        COL2_W - Emu(160000), Emu(200000), size=13, bold=True, color=CDATA_DARK, wrap=False)

    right_items = [
        ('Discovery bottleneck', 'SEO + LLM presence content'),
        ('Evaluation drop-off', 'Comparison content + quickstart fix'),
        ('Activation friction', 'DX improvements + error messages'),
    ]
    dy = TOP + Emu(360000)
    for signal, action in right_items:
        box(s, signal, COL2_X, dy, COL2_W, Emu(200000),
            size=14, bold=True, color=CDATA_DARK, wrap=False)
        box(s, f'→ {action}', COL2_X, dy + Emu(200000), COL2_W, Emu(180000),
            size=13, color=CDATA_MID, wrap=True)
        dy += Emu(480000)

    # Bottom: product conversation track
    shape(s, ML, Emu(3250000), CW, Emu(600000), fill=CDATA_CARD, rounded=False)
    box(s, 'Product track', ML + Emu(200000), Emu(3330000), Emu(1900000), Emu(280000),
        size=14, bold=True, color=CDATA_DARK, wrap=False)
    box(s, 'Take the 47-issue audit + baseline data to the product team — before/after framing, not a wish list.',
        ML + Emu(2100000), Emu(3330000), CW - Emu(2300000), Emu(400000),
        size=14, color=CDATA_MID, wrap=True)
    set_notes(s, SCRIPTS[38])


def slide_38_days_6190(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Days 61–90')
    box(s, '60 days of evidence earns us the right to build this.',
        ML, Emu(620000), CW, Emu(500000),
        size=32, bold=True, color=CDATA_DARK, wrap=True)

    three_cards(s, [
        {
            'title': 'Double Down on What Moved',
            'body':  (
                'Content expands where signal is strongest.\n\n'
                'Comparison piece worked → SAP, NetSuite, Workday quickstarts. '
                'Enterprise sources Composio can\'t touch.\n\n'
                'Community seeding worked → conference talk + OSS working demo.\n\n'
                'We don\'t guess — we expand on proof.'
            ),
        },
        {
            'title': 'Sales Enablement',
            'body':  (
                'Enterprise deals stall when the AE can\'t answer the technical question.\n\n'
                'Battlecard: Connect AI vs Composio — one honest answer per objection.\n\n'
                '"When an engineer asks about Connect AI" — a one-page technical guide for the field.\n\n'
                'PMM ships this. DevRel tests it with community.'
            ),
        },
        {
            'title': 'Day 90: Re-run the Baseline',
            'body':  (
                'Same 5 LLM queries from Day 1.\n\n'
                'Day 1:   0 / 5 mentioned CData.\n'
                'Target:  2–3 / 5.\n\n'
                'This is the number we\'re accountable to. Not traffic, not MQLs. '
                'The metric that proves positioning actually moved.'
            ),
        },
    ], top=Emu(1350000), card_h=Emu(4900000), animate=False)
    set_notes(s, SCRIPTS[39])


def slide_39_goals(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Goals')
    box(s, 'Goals framework.',
        ML, Emu(620000), CW, Emu(700000),
        size=40, bold=True, color=CDATA_DARK, wrap=True)

    rows = [
        ('LLM discoverability (5 key queries)',        '0/5 — current baseline', '2–3/5'),
        ('Time-to-first-API-call',                     '~15 min, 6 friction pts', 'Under 15 min, clean'),
        ('Developer funnel drop-off (by stage)',        'Established at Day 30',  '20%+ improvement at highest-loss stage'),
        ('Keyword cluster rankings (5 clusters)',       'Mapped at Day 30',        'Movement on 3–5 clusters'),
        ('Community-influenced pipeline',               'Not attributed',          'Attribution established + first tagged deal'),
    ]

    header_y = Emu(1700000)
    row_h    = Emu(820000)

    shape(s, ML, header_y, CW, Emu(380000), fill=CDATA_DARK, rounded=False)
    for off, hdr in [(0, 'Metric'), (Emu(5400000), 'Day 30 baseline'), (Emu(8600000), 'Day 90 target')]:
        box(s, hdr, ML + off + Emu(180000), header_y + Emu(80000), Emu(4000000), Emu(240000),
            size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=False)

    for i, (metric, baseline, target) in enumerate(rows):
        row_y = header_y + Emu(380000) + i * row_h
        fill  = CDATA_CARD if i % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF5)
        shape(s, ML, row_y, CW, row_h - Emu(30000), fill=fill, rounded=False)
        box(s, metric,   ML + Emu(180000),              row_y + Emu(140000), Emu(5000000), Emu(560000),
            size=13, color=CDATA_DARK, wrap=True)
        box(s, baseline, ML + Emu(5400000) + Emu(180000), row_y + Emu(140000), Emu(2900000), Emu(560000),
            size=13, color=CDATA_MID, wrap=True)
        box(s, target,   ML + Emu(8600000) + Emu(180000), row_y + Emu(140000), Emu(2000000), Emu(560000),
            size=13, color=CDATA_DARK, wrap=True)

    set_notes(s, SCRIPTS[40])


def slide_40_pmm_devrel(prs):
    s = blank(prs)
    set_bg(s, CDATA_BG)
    add_cdata_logo(s, 'dark')
    overline(s, 'Team Structure')
    box(s, 'PMM and DevRel.',
        ML, Emu(620000), CW, Emu(700000),
        size=44, bold=True, color=CDATA_DARK, wrap=True)

    two_cards(s, [
        {
            'title': 'DevRel earns trust',
            'body':  (
                'Discord, GitHub, conference talks, working demos.\n\n'
                'Builds presence where the developer actually is.\n\n'
                'Brings back what developers are asking and what they\'re struggling with.'
            ),
        },
        {
            'title': 'PMM turns trust into pipeline',
            'body':  (
                'Positioning, comparison content, attribution, Sales enablement.\n\n'
                'Translates developer signal into commercial messaging.\n\n'
                'Builds the content and funnel infrastructure that connects developer adoption to enterprise deals.'
            ),
        },
    ], top=Emu(1800000), card_h=Emu(3200000), animate=False)

    box(s, 'These are different jobs. They fail when blurred. Interface: weekly sync, shared funnel definition.',
        ML, Emu(5300000), CW, Emu(500000),
        size=14, color=CDATA_GRAY, wrap=True)
    set_notes(s, SCRIPTS[41])


def slide_41_content_instinct(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, 'What developer content actually works:',
        ML, Emu(1400000), CW, Emu(600000),
        size=22, color=CDATA_YELLOW, wrap=True)

    box(s, "it's specific, it's runnable, and it's honest about failure modes.",
        ML, Emu(2150000), CW, Emu(1200000),
        size=38, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)

    shape(s, ML, Emu(3600000), CW, Emu(40000), fill=CDATA_DIVIDER, rounded=False)

    box(s, '"Seamless enterprise connectivity" without a code sample gets ignored.',
        ML, Emu(3800000), CW, Emu(500000),
        size=18, color=CDATA_GRAY, wrap=True)

    box(s, '"Here\'s what happens when your PAT expires mid-query and how to handle it" gets bookmarked.',
        ML, Emu(4500000), CW, Emu(600000),
        size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)
    set_notes(s, SCRIPTS[42])


# ── CLOSE ─────────────────────────────────────────────────────────────────────

def slide_42_artifacts(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, "What I'm leaving behind.", ML, Emu(600000), CW, Emu(700000),
        size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)

    artifacts = [
        ('Working demo', 'github.com/JoeKarlsson/connect-ai-agent-demo — clone it and run it. The agent queries live enterprise data.'),
        ('47-issue DX audit', 'docs/dx-audit.md — root causes, recommended fixes, organized by category and impact.'),
        ('Composio competitive brief', 'SQL vs. action-based, enterprise data coverage, pricing comparison, CData\'s defensible position.'),
        ('LLM baseline data', '5 key queries run May 8, 2026. CData 0/5. Composio 5/5. The before numbers.'),
    ]

    dy = Emu(1700000)
    for title, desc in artifacts:
        shape(s, ML, dy, Emu(90000), Emu(90000), fill=CDATA_YELLOW, rounded=False)
        box(s, title, ML + Emu(260000), dy, Emu(3400000), Emu(380000),
            size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), wrap=False)
        box(s, desc, ML + Emu(260000), dy + Emu(400000), Emu(7200000), Emu(380000),
            size=14, color=CDATA_GRAY, wrap=True)
        dy += Emu(900000)

    # QR code
    try:
        from PIL import Image as _I
        qw, qh = _I.open(_REPO_QR).size
        qh_emu = Emu(1400000)
        qw_emu = int(qh_emu * qw / qh)
        s.shapes.add_picture(_REPO_QR,
            W - ML - qw_emu - Emu(200000), Emu(2000000),
            qw_emu, qh_emu)
    except Exception:
        pass

    set_notes(s, SCRIPTS[43])


def slide_43_close(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    lines = [
        ("There's a developer right now who needs to connect their AI agent to Salesforce and NetSuite in a single query.", 20, CDATA_GRAY),
        ("They're going to Google it, ask Claude, check GitHub.", 22, RGBColor(0xFF, 0xFF, 0xFF)),
        ("Right now, Composio comes up. CData doesn't.", 24, CDATA_YELLOW),
        ("That's not a product problem — CData has the product.", 22, RGBColor(0xFF, 0xFF, 0xFF)),
        ("Getting there is a messaging, content, and distribution problem.", 22, RGBColor(0xFF, 0xFF, 0xFF)),
        ("That's exactly what I do.", 28, CDATA_YELLOW),
    ]

    dy = Emu(1000000)
    for text, size, color in lines:
        h = Emu(size * 12700 * 2)
        box(s, text, ML, dy, CW, h, size=size, color=color, wrap=True)
        dy += h + Emu(120000)

    set_notes(s, SCRIPTS[44])


def slide_44_questions(prs):
    s = blank(prs)
    set_bg(s, CDATA_DARK)
    add_cdata_logo(s, 'light')

    box(s, 'Questions.', 0, Emu(2400000), W, Emu(1400000),
        size=88, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    box(s, 'joekarlsson.com', 0, Emu(4200000), W, Emu(500000),
        size=16, color=CDATA_GRAY, align=PP_ALIGN.CENTER,
        url='https://joekarlsson.com')
    set_notes(s, SCRIPTS[45])


# ── Notes injection (post-import into Keynote) ────────────────────────────────

def inject_keynote_notes(doc_name_fragment='cdata-interview'):
    """Re-inject SCRIPTS into an open Keynote document via AppleScript.

    postprocess_pptx strips all notesSlide XML to prevent Keynote Code=2 errors.
    Run this AFTER opening the pptx in Keynote to restore speaker notes.

    Usage:
        1. Open cdata-interview.pptx in Keynote
        2. .venv/bin/python3 presentation/build_slides.py --inject-notes
    """
    import subprocess
    lines = ['tell application "Keynote"']
    lines.append('    set targetDoc to missing value')
    lines.append('    repeat with d in documents')
    lines.append(f'        if name of d contains "{doc_name_fragment}" then')
    lines.append('            set targetDoc to d')
    lines.append('            exit repeat')
    lines.append('        end if')
    lines.append('    end repeat')
    lines.append('    if targetDoc is missing value then')
    lines.append(f'        error "Could not find Keynote document containing \\"{doc_name_fragment}\\""')
    lines.append('    end if')
    for slide_num, script_text in SCRIPTS.items():
        escaped = script_text.replace('\\', '\\\\').replace('"', '\\"').replace('\n', '\\n')
        lines.append(f'    set presenter notes of slide {slide_num} of targetDoc to "{escaped}"')
    lines.append('end tell')
    result = subprocess.run(['osascript', '-e', '\n'.join(lines)], capture_output=True, text=True)
    if result.returncode != 0:
        print(f'  [notes] AppleScript error: {result.stderr.strip()}')
        print('  [notes] Make sure cdata-interview.pptx is open in Keynote first.')
    else:
        print(f'  [notes] Injected notes for {len(SCRIPTS)} slides into Keynote.')


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    check_fonts()
    prs = new_prs()

    slide_01_title(prs)
    slide_02_before_strategy(prs)
    slide_03_demo(prs)
    slide_04_demo_result(prs)
    slide_05_built_and_broke(prs)

    slide_06_developer_who_needs(prs)
    slide_07_cross_source_pain(prs)
    slide_08_enterprise_data(prs)
    slide_09_cdata_position(prs)
    slide_10_not_reached(prs)
    slide_11_three_audiences(prs)
    slide_12_llm_agents(prs)
    slide_13_ai_developers(prs)
    slide_14_enterprise_buyers(prs)
    slide_15_different_approach(prs)

    slide_16_whats_working(prs)
    slide_17_issue_table(prs)
    slide_18_not_all_47(prs)
    slide_19_ai_doc_standards(prs)
    slide_20_finding1_intro(prs)
    slide_21_screenshot(prs)
    slide_21_5_queries(prs)
    slide_22_llms_where_dev_starts(prs)
    slide_23_finding2_no_agent_path(prs)
    slide_24_same_seven_steps(prs)
    slide_25_cta_not_docs(prs)
    slide_26_finding3_orphaned(prs)
    slide_27_no_path_to_docs(prs)
    slide_28_wrong_place(prs)
    slide_29_the_pattern(prs)
    slide_30_two_buckets(prs)
    slide_31_pmm_product_bridge(prs)

    slide_32_90_days(prs)
    slide_33_days_130_intro(prs)
    slide_34_three_workstreams(prs)
    slide_35_immediate_fixes(prs)
    slide_36_entry_point(prs)
    slide_37_days_3160(prs)
    slide_38_days_6190(prs)
    slide_39_goals(prs)
    slide_40_pmm_devrel(prs)
    slide_41_content_instinct(prs)

    slide_42_artifacts(prs)
    slide_43_close(prs)
    slide_44_questions(prs)

    out = os.path.join(HERE, 'cdata-interview.pptx')
    prs.save(out)
    postprocess_pptx(out)
    print(f'Saved {len(prs.slides)} slides → {out}')


if __name__ == '__main__':
    if '--inject-notes' in sys.argv:
        inject_keynote_notes()
    else:
        build()
