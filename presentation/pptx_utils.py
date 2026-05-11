"""
env0 2026 brand utilities for python-pptx presentations.

Usage in any build_slides.py:

    import sys, os
    HERE = os.path.dirname(os.path.abspath(__file__))
    sys.path.insert(0, os.path.join(HERE, '..'))
    from pptx_utils import *

Everything is then available: colors, layout constants, helpers, post-processing.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, uuid

# ── Asset paths ───────────────────────────────────────────────────────────────
# Template assets live in presentations/template-assets/ (shared across all talks)
ASSETS   = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'template-assets')
LOGO_IMG      = os.path.join(ASSETS, 'template-logo.png')
LOGO_IMG_LIME = os.path.join(ASSETS, 'template-logo-lime.png')
SAGE_DIA = os.path.join(ASSETS, 'template-diamonds-sage.png')
LIME_DIA = os.path.join(ASSETS, 'template-diamonds-lime.png')
STRIPES  = os.path.join(ASSETS, 'template-stripes.png')

# ── Colors ─────────────────────────────────────────────────────────────────────
BG      = RGBColor(0xF4, 0xF5, 0xF0)  # warm off-white — content slide background
BG_HERO = RGBColor(0xFF, 0xFF, 0xFF)  # white — hero/title slide background
BG_LIME = RGBColor(0xD2, 0xF7, 0x40)  # lime — mandate/divider background
GREEN   = RGBColor(0xD2, 0xF7, 0x40)  # lime — fills/shapes ONLY, never text on light bg
TEAL    = RGBColor(0x03, 0x99, 0x8B)  # teal — overline labels, card titles, accent text
DARK    = RGBColor(0x00, 0x00, 0x00)  # black — h1 titles, headings
GRAY    = RGBColor(0x80, 0x80, 0x80)  # gray — subtitles, secondary text
LGRAY   = RGBColor(0x4A, 0x4A, 0x4A)  # dark gray — body text on light bg
MID     = LGRAY
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)  # white — text on dark/lime fills
CARD    = RGBColor(0xEC, 0xED, 0xE8)  # card background
DIVIDER = RGBColor(0xD8, 0xDB, 0xCF)  # divider lines / horizontal rules
GHBG    = RGBColor(0xE8, 0xF5, 0xDC)  # light lime tint — highlight box background
CODEBG  = RGBColor(0xEC, 0xED, 0xE8)  # code block background
RED     = RGBColor(0xEF, 0x53, 0x7B)  # coral — error / negative highlights
AMBER   = RGBColor(0xFF, 0xC6, 0x0B)  # amber — warning highlights
ERRRED  = RED

# ── Typography ─────────────────────────────────────────────────────────────────
FONT    = 'Geist'
FONT_XB = 'Geist ExtraBold'  # user switched all bold Geist text to this weight
MONO    = 'Courier New'

# ── Canvas — 13.333" × 7.5" ───────────────────────────────────────────────────
W = Emu(12192000)
H = Emu(6858000)

# ── Layout grid ───────────────────────────────────────────────────────────────
TXT_L = Emu(685628)    # left margin (0.75")
TXT_T = Emu(544550)    # h1 top (0.60")
TXT_W = Emu(5803500)   # left text zone — clears diamond pattern

ML  = TXT_L
MR  = TXT_L
CW  = W - ML - MR      # full content width = 10,820,744 EMU

H1_H        = Emu(1542000)  # h1 box height — user layout: ~42.8mm
SUB_Y       = Emu(2140000)  # subtitle top — user layout: ~59.4mm
CONTENT_T   = Emu(2800000)  # content start without subtitle
CONT_S      = Emu(3520000)  # content start after subtitle
SAFE_BOTTOM = Emu(5700000)  # hard lower limit for content

# ── Logo — upper-right corner ─────────────────────────────────────────────────
LOGO_W = Emu(2422145)
LOGO_H = Emu(318945)
LOGO_L = Emu(12192000) - Emu(685628) - LOGO_W   # right-aligned with margin
LOGO_T = Emu(291000)                              # vertically centered with overline (overline center = 450000)

# ── Template decorative asset positions ───────────────────────────────────────
SAGE_L, SAGE_T = Emu(5501723), Emu(0)
SAGE_W, SAGE_H = Emu(6687229), Emu(6857999)

LIME_L, LIME_T = Emu(5501733), Emu(0)
LIME_W, LIME_H = Emu(6687219), Emu(6857999)

LP_L, LP_T = Emu(8067681), Emu(5400)    # lime right panel (hero slides)
LP_W, LP_H = Emu(4121400), Emu(6858000)

STR_L, STR_T = Emu(6888075), Emu(0)     # chevron stripes (hero slides)
STR_W, STR_H = Emu(5300881), Emu(6857999)

# ── Namespaces (for animation XML) ────────────────────────────────────────────
_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


# ── Font preflight ────────────────────────────────────────────────────────────

def check_fonts():
    """Check that required fonts are installed. Prints install instructions if any are missing.

    Call at the top of build() so speakers get an actionable error before the
    build runs, rather than a silent font-substitution warning in Keynote.
    Returns True if all fonts present, False if any are missing (build continues either way).
    """
    user_fonts = os.path.expanduser('~/Library/Fonts')
    required = {
        'Geist-Regular.otf':     os.path.join(user_fonts, 'Geist-Regular.otf'),
        'Geist-Bold.otf':        os.path.join(user_fonts, 'Geist-Bold.otf'),
        'Geist-ExtraBold.otf':   os.path.join(user_fonts, 'Geist-ExtraBold.otf'),
        'DMSans[opsz,wght].ttf': os.path.join(user_fonts, 'DMSans[opsz,wght].ttf'),
    }
    missing = [name for name, path in required.items() if not os.path.exists(path)]
    if not missing:
        return True
    print(f'[fonts] Missing: {", ".join(missing)}')
    print('[fonts] Install with:')
    print('  brew install --cask font-geist font-dm-sans')
    print('  mkdir -p ~/Library/Fonts/disabled')
    print('  mv ~/Library/Fonts/Geist[wght].ttf ~/Library/Fonts/Geist-Italic[wght].ttf \\')
    print('     ~/Library/Fonts/disabled/ 2>/dev/null || true')
    print('[fonts] Deck will build but Keynote will show font substitution warnings.')
    return False


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    """Create a new Presentation with env0 canvas dimensions."""
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    """Add a fully blank slide (layout index 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=None):
    """Set slide background. Default = BG (warm off-white)."""
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color if color is not None else BG

def add_logo(slide):
    """Upper-right env0 logo (dark) — always call LAST so it renders on top."""
    slide.shapes.add_picture(LOGO_IMG, LOGO_L, LOGO_T, LOGO_W, LOGO_H)

def add_logo_lime(slide):
    """Upper-right env0 logo (lime) — for dark/meme backgrounds."""
    slide.shapes.add_picture(LOGO_IMG_LIME, LOGO_L, LOGO_T, LOGO_W, LOGO_H)

def add_sage_diamonds(slide):
    """Sage diamond pattern — call FIRST (behind all content). Use on content slides."""
    slide.shapes.add_picture(SAGE_DIA, SAGE_L, SAGE_T, SAGE_W, SAGE_H)

def add_lime_diamonds(slide):
    """Lime diamond pattern — call FIRST. Use on divider/stat slides."""
    slide.shapes.add_picture(LIME_DIA, LIME_L, LIME_T, LIME_W, LIME_H)

def add_hero_stripes(slide):
    """Chevron stripe overlay — call after lime panel shape on hero slides."""
    slide.shapes.add_picture(STRIPES, STR_L, STR_T, STR_W, STR_H)

def box(slide, text, l, t, w, h,
        size=14, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT,
        wrap=True, font=None, url=None):
    """Add a single-run text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    f  = r.font
    effective_font = font or FONT
    if bold and effective_font == FONT:
        effective_font = FONT_XB  # user changed all bold Geist → Geist ExtraBold
        bold = False
    f.name      = effective_font
    f.size      = Pt(size)
    f.bold      = bold
    f.italic    = italic
    f.color.rgb = color
    if url:
        r.hyperlink.address = url
        f.underline = False
    return tb

def mbox(slide, paras, l, t, w, h):
    """Add a multi-paragraph text box.

    paras: list of para dicts. Two forms:

    Single-run (original):
        {'text', 'size', 'bold', 'italic', 'color', 'font', 'align', 'space_after', 'url'}

    Multi-run (use when you need inline hyperlinks or mixed styles in one line):
        {'runs': [{'text', 'size', 'bold', 'italic', 'color', 'font', 'url'}, ...],
         'align', 'space_after'}
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pd.get('align', PP_ALIGN.LEFT)
        sp = pd.get('space_after')
        if sp:
            p.space_after = Pt(sp)
        runs = pd.get('runs')
        if runs:
            for rd in runs:
                r = p.add_run()
                r.text      = rd.get('text', '')
                f           = r.font
                f.name      = rd.get('font', FONT)
                f.size      = Pt(rd.get('size', 14))
                f.bold      = rd.get('bold', False)
                f.italic    = rd.get('italic', False)
                f.color.rgb = rd.get('color', DARK)
                if rd.get('url'):
                    r.hyperlink.address = rd['url']
        else:
            r = p.add_run()
            r.text      = pd.get('text', '')
            f           = r.font
            f.name      = pd.get('font', FONT)
            f.size      = Pt(pd.get('size', 14))
            f.bold      = pd.get('bold', False)
            f.italic    = pd.get('italic', False)
            f.color.rgb = pd.get('color', DARK)
            if pd.get('url'):
                r.hyperlink.address = pd['url']
    return tb

def shape(slide, l, t, w, h, fill=CARD, rounded=True, border=None):
    """Add a filled background rectangle — DECORATION ONLY, never add text to it.

    Keynote requires text to live in proper text boxes (txBox="1"). Shapes
    created here do NOT have txBox="1", so accessing .text_frame and adding
    text will produce slides that Keynote rejects with "cannot import" errors.

    To show text over a card: call shape() for the background, then call box()
    or mbox() with matching coordinates to overlay text on top.
    """
    stype = 5 if rounded else 1
    sh = slide.shapes.add_shape(stype, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh

def set_notes(slide, text):
    """Add speaker notes text to a slide."""
    slide.notes_slide.notes_text_frame.text = text

def source_label(slide, text, url=None, top=None):
    """Small 9pt gray source citation near the bottom of the slide."""
    t = top or Emu(6300000)
    box(slide, f'SOURCE  {text}', ML, t,
        CW, Emu(260000),
        size=9, color=GRAY, wrap=False, url=url)


# ── Layout components ─────────────────────────────────────────────────────────

def overline(slide, text, left=None):
    """Small TEAL uppercase category label — DM Sans 11pt Regular (user switched from Bold).
    Sits above the h1 at y=340,000 EMU."""
    box(slide, text.upper(), left or ML, Emu(340000), TXT_W, Emu(220000),
        size=11, bold=False, color=TEAL, font='DM Sans')

def h1(slide, text, top=None, size=44, width=None):
    """Main slide title — Geist Bold 44pt black.
    Default position: y=TXT_T, height=H1_H (fits 3 lines)."""
    box(slide, text, ML, top or TXT_T, width or TXT_W, H1_H,
        size=size, bold=True, color=DARK, wrap=True)

def sub(slide, text, top=None, size=22, width=None, color=GRAY):
    """Subtitle — Geist Regular 22pt gray.
    Default position: y=SUB_Y."""
    box(slide, text, ML, top or SUB_Y, width or TXT_W, Emu(700000),
        size=size, color=color, wrap=True)

def highlight_box(slide, text, top, width=None, height=None):
    """Light-lime callout box with a left teal accent bar.
    Good for quotes, key insights, or command output.
    Height auto-sizes to fit text if not explicitly provided."""
    import math as _math, textwrap as _tw
    l = ML
    w = width or CW
    if height is None:
        txt_w   = w - Emu(260000)
        cpl     = max(20, int(txt_w / 12700 / 7))   # chars/line at ~7pt avg char width
        n_lines = max(1, len(_tw.wrap(text, cpl)))
        height  = int(n_lines * int(13 * 12700 * 1.5) + Emu(220000))
    bg   = shape(slide, l, top, w, height, fill=GHBG, rounded=False)
    bar  = shape(slide, l, top, Emu(60000), height, fill=GREEN, rounded=False)
    body = box(slide, text, l + Emu(170000), top + Emu(100000),
               w - Emu(260000), height - Emu(170000),
               size=13, color=DARK, wrap=True)
    return [bg.shape_id, bar.shape_id, body.shape_id]

def footer_line(slide, text, top=None):
    """Thin divider + small footnote text near the bottom of the slide."""
    t = top or Emu(6260000)
    shape(slide, ML, t, CW, Emu(22000), fill=DIVIDER, rounded=False)
    box(slide, text, ML, t + Emu(55000), CW, Emu(430000),
        size=11, color=GRAY, wrap=True)

def three_cards(slide, items, top=None, card_h=None, animate=True):
    """Three equal-width cards across the full content width.

    items: list of 3 dicts, each with 'title' and 'body' keys.
    Cards have click-to-appear animation by default (animate=True).
    """
    top    = top    or CONTENT_T
    card_h = card_h or Emu(3080000)
    gap = Emu(185000)
    cw  = (CW - 2 * gap) / 3
    pad = Emu(185000)
    groups = []
    for i, item in enumerate(items):
        l = ML + i * (cw + gap)
        bg  = shape(slide, l, top, cw, card_h, fill=CARD)
        ttl = box(slide, item['title'], l+pad, top+pad, cw-2*pad, Emu(430000),
                  size=17, bold=True, color=TEAL, wrap=True)
        bt  = top + pad + Emu(460000)
        bdy = box(slide, item['body'], l+pad, bt, cw-2*pad,
                  card_h - (bt - top) - pad,
                  size=15, color=LGRAY, wrap=True)
        groups.append([bg.shape_id, ttl.shape_id, bdy.shape_id])
    if animate:
        add_appear_animations(slide, groups)

def two_cards(slide, items, top=None, card_h=None, animate=True):
    """Two equal-width cards. Same API as three_cards."""
    top    = top    or CONTENT_T
    card_h = card_h or Emu(3080000)
    gap = Emu(220000)
    cw  = (CW - gap) / 2
    pad = Emu(185000)
    groups = []
    for i, item in enumerate(items):
        l = ML + i * (cw + gap)
        bg  = shape(slide, l, top, cw, card_h, fill=CARD)
        ttl = box(slide, item['title'], l+pad, top+pad, cw-2*pad, Emu(430000),
                  size=17, bold=True, color=TEAL, wrap=True)
        bt  = top + pad + Emu(460000)
        bdy = box(slide, item['body'], l+pad, bt, cw-2*pad,
                  card_h - (bt - top) - pad,
                  size=15, color=LGRAY, wrap=True)
        groups.append([bg.shape_id, ttl.shape_id, bdy.shape_id])
    if animate:
        add_appear_animations(slide, groups)


# ── Meme slides ──────────────────────────────────────────────────────────────

def meme_slide(prs, media, caption=None, notes=None, caption_top=False, caption_size=36,
               fallback_url=None, quote=None, quote_size=28):
    """Full-slide meme on black background. Accepts local MP4 video paths or GIF/image URLs.

    Local video files (mp4, mov): embedded in the PPTX via add_movie(), poster frame
    extracted with ffmpeg. Plays in Keynote and PowerPoint offline.
    URL strings: downloaded and embedded as a static picture (original behavior).
    fallback_url: if media is a local path that does not exist, use this URL instead.
    Caption is rendered in Impact font, white fill, black stroke — classic meme style.
    """
    import urllib.request as _ur, tempfile as _tf, os as _oss, subprocess as _sp, io as _io

    s = blank(prs)
    set_bg(s, DARK)

    # Fall back to URL when media is None or a local path that doesn't exist
    if fallback_url and (media is None or (isinstance(media, str)
            and not _oss.path.isfile(media)
            and not media.startswith('http'))):
        media = fallback_url

    _is_video = (
        isinstance(media, str)
        and _oss.path.isfile(media)
        and media.lower().endswith(('.mp4', '.mov', '.webm'))
    )

    tmp = None
    _poster_tmp = None
    try:
        if _is_video:
            _poster_tmp = _tf.NamedTemporaryFile(suffix='.jpg', delete=False)
            _poster_tmp.close()
            _sp.run(
                ['ffmpeg', '-i', media, '-vframes', '1', '-q:v', '2', _poster_tmp.name, '-y'],
                capture_output=True
            )
            aspect = 1.0
            try:
                from PIL import Image as _PIL
                with _PIL.open(_poster_tmp.name) as im:
                    w, h = im.size
                    aspect = w / h
            except Exception:
                pass
        else:
            ext = '.' + media.split('?')[0].rsplit('.', 1)[-1]
            if ext not in ('.gif', '.jpg', '.jpeg', '.png', '.webp'):
                ext = '.jpg'
            tmp = _tf.NamedTemporaryFile(suffix=ext, delete=False)
            req = _ur.Request(media, headers={'User-Agent': 'CloudQuery-Slides/1.0'})
            with _ur.urlopen(req) as resp:
                tmp.write(resp.read())
            tmp.close()
            aspect = 1.0
            try:
                from PIL import Image as _PIL
                with _PIL.open(tmp.name) as im:
                    w, h = im.size
                    aspect = w / h
            except Exception:
                pass

        cap_h    = Emu(1200000) if caption else Emu(0)  # room for 2 lines of 36pt Impact
        margin_h = Emu(740000)
        margin_v = Emu(370000)
        if caption_top:
            cap_y   = Emu(309600)            # top of slide (~8.6mm)
            img_top = cap_y + cap_h if caption else margin_v
            avail_h = H - img_top - margin_v
        else:
            avail_h = H - 2 * margin_v - cap_h
            img_top = margin_v
            cap_y   = H - cap_h             # bottom of slide
        avail_w  = W - 2 * margin_h

        if aspect >= avail_w / avail_h:
            pic_w = int(avail_w)
            pic_h = int(avail_w / aspect)
        else:
            pic_h = int(avail_h)
            pic_w = int(avail_h * aspect)

        pic_x = (W - pic_w) // 2
        pic_y = int(img_top + (avail_h - pic_h) // 2)

        if _is_video:
            with open(_poster_tmp.name, 'rb') as pf:
                poster = _io.BytesIO(pf.read())
            s.shapes.add_movie(media, pic_x, pic_y, pic_w, pic_h,
                               poster_frame_image=poster, mime_type='video/mp4')
        else:
            s.shapes.add_picture(tmp.name, pic_x, pic_y, pic_w, pic_h)

        if caption:
            # Classic meme caption: Impact, white fill, black outline stroke, ALL CAPS
            tb  = s.shapes.add_textbox(ML, cap_y, CW, cap_h)
            tf  = tb.text_frame
            tf.word_wrap = True
            par = tf.paragraphs[0]
            par.alignment = PP_ALIGN.CENTER
            run = par.add_run()
            run.text = caption.upper()
            rPr = run._r.get_or_add_rPr()
            rPr.set('sz', str(caption_size * 100))  # caption_size pt
            rPr.set('b', '0')
            # 1) Stroke/outline — must precede fill in CT_RPr schema order
            ln = etree.SubElement(rPr, qn('a:ln'))
            ln.set('w', '19050')    # 1.5pt = 1.5 × 12700 EMU
            lf = etree.SubElement(ln, qn('a:solidFill'))
            lc = etree.SubElement(lf, qn('a:srgbClr'))
            lc.set('val', '000000')
            # 2) Fill
            sf = etree.SubElement(rPr, qn('a:solidFill'))
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', 'FFFFFF')
            # 3) Font
            lt = etree.SubElement(rPr, qn('a:latin'))
            lt.set('typeface', 'Impact')

        if quote:
            # Subtitle-style overlay: yellow Impact text with black stroke, bottom of video frame
            q_h   = Emu(600000)
            q_pad = Emu(80000)
            q_top = pic_y + pic_h - q_h - q_pad
            q_tb  = s.shapes.add_textbox(pic_x, q_top, pic_w, q_h)
            q_tf  = q_tb.text_frame
            q_tf.word_wrap = True
            q_par = q_tf.paragraphs[0]
            q_par.alignment = PP_ALIGN.CENTER
            q_run = q_par.add_run()
            q_run.text = quote
            q_rPr = q_run._r.get_or_add_rPr()
            q_rPr.set('sz', str(quote_size * 100))
            q_rPr.set('b', '1')
            q_ln = etree.SubElement(q_rPr, qn('a:ln'))
            q_ln.set('w', '25400')   # 2pt stroke for legibility over video
            q_lf = etree.SubElement(q_ln, qn('a:solidFill'))
            q_lc = etree.SubElement(q_lf, qn('a:srgbClr'))
            q_lc.set('val', '000000')
            q_sf = etree.SubElement(q_rPr, qn('a:solidFill'))
            q_sc = etree.SubElement(q_sf, qn('a:srgbClr'))
            q_sc.set('val', 'FFE000')   # classic subtitle yellow
            q_lt = etree.SubElement(q_rPr, qn('a:latin'))
            q_lt.set('typeface', 'Impact')
    finally:
        if tmp:
            try:
                _oss.unlink(tmp.name)
            except Exception:
                pass
        if _poster_tmp:
            try:
                _oss.unlink(_poster_tmp.name)
            except Exception:
                pass

    add_logo_lime(s)
    if notes:
        set_notes(s, notes)
    return s


# ── Click-to-appear animations ────────────────────────────────────────────────

def add_appear_animations(slide, groups):
    """Add click-to-appear animations to groups of shape IDs.

    groups: list of lists of shape IDs. Each inner list appears together on one click.
    """
    if not groups:
        return
    cid = [200]
    def nid():
        cid[0] += 1
        return str(cid[0])

    def effect_par(spid, grp_id, node_type):
        eid, sid = nid(), nid()
        return (
            f'<p:par xmlns:p="{_P}">'
            f'<p:cTn id="{eid}" presetID="1" presetClass="entr" presetSubtype="0"'
            f' fill="hold" grpId="{grp_id}" nodeType="{node_type}">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:set>'
            f'<p:cBhvr><p:cTn id="{sid}" dur="1" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr>'
            f'<p:to><p:strVal val="visible"/></p:to>'
            f'</p:set></p:childTnLst></p:cTn></p:par>'
        )

    click_pars = []
    bld_entries = []
    for grp_idx, id_list in enumerate(groups):
        outer_id = nid()
        effects = []
        for k, spid in enumerate(id_list):
            node_type = 'clickEffect' if k == 0 else 'withEffect'
            effects.append(effect_par(spid, grp_idx, node_type))
            bld_entries.append(
                f'<p:bldP xmlns:p="{_P}" spid="{spid}"'
                f' grpId="{grp_idx}" uiExpand="1" build="allAtOnce"/>'
            )
        click_pars.append(
            f'<p:par xmlns:p="{_P}">'
            f'<p:cTn id="{outer_id}" fill="hold">'
            f'<p:stCondLst><p:cond evt="onBegin" delay="indefinite"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(effects)}</p:childTnLst>'
            f'</p:cTn></p:par>'
        )

    root_id, seq_id = nid(), nid()
    timing_xml = (
        f'<p:timing xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}">'
        f'<p:tnLst><p:par>'
        f'<p:cTn id="{root_id}" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">'
        f'<p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="{seq_id}" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(click_pars)}</p:childTnLst>'
        f'</p:cTn>'
        f'<p:prevCondLst><p:cond evt="onBegin" delay="indefinite"/></p:prevCondLst>'
        f'</p:seq>'
        f'</p:childTnLst></p:cTn>'
        f'</p:par></p:tnLst>'
        f'<p:bldLst>{"".join(bld_entries)}</p:bldLst>'
        f'</p:timing>'
    )
    slide._element.append(etree.fromstring(timing_xml))


# ── Post-processing ───────────────────────────────────────────────────────────

def fix_pptx_for_keynote(path):
    """Patch OOXML for Apple Keynote compatibility.

    Fixes three issues:
    1. Sets sldSz type="custom" so Keynote respects the 13.333"x7.5" canvas.
    2. Removes dangling notesMaster references when the notes master XML file
       doesn't exist in the zip — python-pptx adds the relationship but skips
       the actual file, and Keynote hard-rejects decks with unresolvable refs.
    3. Replaces the default python-pptx font scheme (full of Windows-only
       script fonts) with a minimal Geist-only scheme so Keynote doesn't show
       a "replace fonts" dialog on every open.
    """
    import zipfile, re
    with zipfile.ZipFile(path, 'r') as zin:
        files = {item.filename: zin.read(item.filename) for item in zin.infolist()}
        infos = {item.filename: item for item in zin.infolist()}

    # Fix sldSz type="custom"
    prs_xml = files['ppt/presentation.xml'].decode('utf-8')
    prs_xml = re.sub(r'(<p:sldSz\b[^>]*)\btype="[^"]*"', r'\1type="custom"', prs_xml)
    m2 = re.search(r'<p:sldSz[^>]*>', prs_xml)
    if m2 and 'type=' not in m2.group():
        prs_xml = re.sub(r'(<p:sldSz\b)', r'\1 type="custom"', prs_xml)

    # Remove dangling notesMaster refs when the actual XML file is absent.
    # python-pptx registers a notesMaster relationship but never writes the file;
    # Keynote rejects the entire deck if it finds an unresolvable relationship.
    if 'ppt/notesMasters/notesMaster1.xml' not in files:
        prs_xml = re.sub(
            r'<p:notesMasterIdLst>.*?</p:notesMasterIdLst>', '',
            prs_xml, flags=re.DOTALL
        )
        rels_xml = files['ppt/_rels/presentation.xml.rels'].decode('utf-8')
        rels_xml = re.sub(
            r'\s*<Relationship[^>]+notesMaster[^>]*/>',
            '',
            rels_xml
        )
        files['ppt/_rels/presentation.xml.rels'] = rels_xml.encode('utf-8')

    files['ppt/presentation.xml'] = prs_xml.encode('utf-8')

    # Replace the theme's font scheme with a minimal Geist-only scheme.
    # python-pptx's default theme contains dozens of Windows-specific script
    # fonts (MS PGothic, Mongolian Baiti, etc.) that don't exist on macOS;
    # Keynote shows a "replace fonts" dialog for every one of them.
    # All our slide text uses explicit typeface="Geist" or "DM Sans", so the
    # theme font scheme is only used for the slide master placeholders — safe
    # to replace with Geist.
    clean_font_scheme = (
        '<a:fontScheme name="env0">'
        '<a:majorFont><a:latin typeface="Geist"/><a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
        '<a:minorFont><a:latin typeface="Geist"/><a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
        '</a:fontScheme>'
    )
    for theme_key in [k for k in files if k.startswith('ppt/theme/') and k.endswith('.xml')]:
        theme_xml = files[theme_key].decode('utf-8')
        theme_xml = re.sub(
            r'<a:fontScheme\b.*?</a:fontScheme>',
            clean_font_scheme,
            theme_xml, flags=re.DOTALL
        )
        files[theme_key] = theme_xml.encode('utf-8')

    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for fname, data in files.items():
            zout.writestr(infos[fname], data)


def embed_fonts(path):
    """Embed Geist Regular, Geist Bold, and DM Sans into the PPTX zip.

    Fonts are read from ~/Library/Fonts (macOS). If a font file is missing,
    that variant is skipped and a warning is printed — the deck still works,
    viewers just need the font installed locally.

    OOXML requires:
    - Font data in ppt/fonts/*.fntdata (first 32 bytes XOR-obfuscated with GUID key)
    - Relationships in ppt/_rels/presentation.xml.rels
    - <p:embeddedFontLst> entries in ppt/presentation.xml
    - <Default Extension="fntdata"> in [Content_Types].xml  ← easy to miss
    """
    import zipfile, re

    user_fonts = os.path.expanduser('~/Library/Fonts')
    font_files = {
        'Geist': {
            'regular': os.path.join(user_fonts, 'Geist-Regular.otf'),
            'bold':    os.path.join(user_fonts, 'Geist-Bold.otf'),
        },
        'DM Sans': {
            'regular': os.path.join(user_fonts, 'DMSans[opsz,wght].ttf'),
        },
    }

    def _obfuscate(font_data, guid_str):
        key = bytes.fromhex(guid_str.strip('{}').replace('-', ''))
        data = bytearray(font_data)
        for i in range(min(32, len(data))):
            data[i] ^= key[i % 16]
        return bytes(data)

    with zipfile.ZipFile(path, 'r') as zin:
        files = {item.filename: zin.read(item.filename) for item in zin.infolist()}
        infos = {item.filename: item for item in zin.infolist()}

    rels_xml = files['ppt/_rels/presentation.xml.rels'].decode('utf-8')
    prs_xml  = files['ppt/presentation.xml'].decode('utf-8')

    rids = [int(x) for x in re.findall(r'rId(\d+)', rels_xml)]
    next_rid = max(rids) + 1 if rids else 30

    font_num   = 1
    new_rels   = []
    font_nodes = []

    variant_tags = {'regular': 'regular', 'bold': 'bold',
                    'italic': 'italic', 'boldItalic': 'boldItalic'}

    for typeface, variants in font_files.items():
        var_parts = []
        for variant_key, fpath in variants.items():
            if not os.path.exists(fpath):
                print(f'  [embed] SKIP missing: {fpath}')
                continue
            with open(fpath, 'rb') as f:
                raw = f.read()
            g = '{' + str(uuid.uuid4()).upper() + '}'
            obf = _obfuscate(raw, g)
            fname = f'ppt/fonts/font{font_num}.fntdata'
            files[fname] = obf
            rid = f'rId{next_rid}'
            new_rels.append(
                f'<Relationship Id="{rid}" '
                f'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
                f'relationships/font" Target="fonts/font{font_num}.fntdata"/>'
            )
            tag = variant_tags[variant_key]
            var_parts.append(f'<p:{tag} r:id="{rid}" uniqueId="{g}"/>')
            font_num  += 1
            next_rid  += 1
        if var_parts:
            node = (f'<p:embeddedFont>'
                    f'<p:font typeface="{typeface}" charset="0" pitchFamily="34"/>'
                    + ''.join(var_parts) +
                    f'</p:embeddedFont>')
            font_nodes.append(node)
            print(f'  [embed] {typeface}: {list(variants.keys())}')

    if not font_nodes:
        print('  [embed] no fonts embedded — check ~/Library/Fonts for Geist-Regular.otf etc.')
        return

    # Register .fntdata content type — without this Office/Keynote silently ignores the files
    ct_xml = files['[Content_Types].xml'].decode('utf-8')
    if 'fntdata' not in ct_xml:
        ct_xml = ct_xml.replace(
            '</Types>',
            '<Default Extension="fntdata" ContentType="application/x-fontdata"/></Types>'
        )
        files['[Content_Types].xml'] = ct_xml.encode('utf-8')

    rels_xml = rels_xml.replace('</Relationships>',
                                '\n'.join(new_rels) + '\n</Relationships>')
    files['ppt/_rels/presentation.xml.rels'] = rels_xml.encode('utf-8')

    embed_block = '<p:embeddedFontLst>' + ''.join(font_nodes) + '</p:embeddedFontLst>'
    if 'embeddedFontLst' not in prs_xml:
        prs_xml = prs_xml.replace('</p:presentation>', embed_block + '</p:presentation>')
    files['ppt/presentation.xml'] = prs_xml.encode('utf-8')

    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for fname, data in files.items():
            if fname in infos:
                zout.writestr(infos[fname], data)
            else:
                zout.writestr(fname, data)


def postprocess_pptx(path):
    """One-pass post-processing: Keynote compatibility fixes + font embedding.

    Supersedes calling fix_pptx_for_keynote() then embed_fonts() separately —
    same work, single zip read/write cycle.
    """
    import zipfile, re

    with zipfile.ZipFile(path, 'r') as zin:
        files = {item.filename: zin.read(item.filename) for item in zin.infolist()}
        infos = {item.filename: item for item in zin.infolist()}

    prs_xml  = files['ppt/presentation.xml'].decode('utf-8')
    rels_xml = files['ppt/_rels/presentation.xml.rels'].decode('utf-8')

    # 1. sldSz type="custom"
    prs_xml = re.sub(r'(<p:sldSz\b[^>]*)\btype="[^"]*"', r'\1type="custom"', prs_xml)
    m = re.search(r'<p:sldSz[^>]*>', prs_xml)
    if m and 'type=' not in m.group():
        prs_xml = re.sub(r'(<p:sldSz\b)', r'\1 type="custom"', prs_xml)

    # 2. Strip all notesSlides and notesMaster — Keynote ignores speaker notes from
    #    PPTX import, and slides that combine a notesSlide ref + an external hyperlink
    #    in the same rels file trigger Compatibility.PptImport Code=2 in Keynote.
    for fname in list(files.keys()):
        if 'notesSlides' in fname or 'notesMasters' in fname:
            del files[fname]
    for slide_rels_key in [k for k in files if k.startswith('ppt/slides/_rels/')]:
        sr = files[slide_rels_key].decode('utf-8')
        sr = re.sub(r'\s*<Relationship[^>]+notesSlide[^>]*/>', '', sr)
        files[slide_rels_key] = sr.encode('utf-8')
    prs_xml  = re.sub(r'<p:notesMasterIdLst>.*?</p:notesMasterIdLst>', '',
                      prs_xml, flags=re.DOTALL)
    rels_xml = re.sub(r'\s*<Relationship[^>]+notesMaster[^>]*/>', '', rels_xml)
    ct_xml = files['[Content_Types].xml'].decode('utf-8')
    ct_xml = re.sub(r'\s*<Override PartName="/ppt/notesSlides[^"]*"[^/]*/>', '', ct_xml)
    ct_xml = re.sub(r'\s*<Override PartName="/ppt/notesMasters[^"]*"[^/]*/>', '', ct_xml)
    files['[Content_Types].xml'] = ct_xml.encode('utf-8')

    # 3. Minimal Geist-only font scheme (strips Windows-only script fonts)
    clean_font_scheme = (
        '<a:fontScheme name="env0">'
        '<a:majorFont><a:latin typeface="Geist"/><a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
        '<a:minorFont><a:latin typeface="Geist"/><a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
        '</a:fontScheme>'
    )
    for theme_key in [k for k in files if k.startswith('ppt/theme/') and k.endswith('.xml')]:
        theme_xml = files[theme_key].decode('utf-8')
        theme_xml = re.sub(r'<a:fontScheme\b.*?</a:fontScheme>', clean_font_scheme,
                           theme_xml, flags=re.DOTALL)
        files[theme_key] = theme_xml.encode('utf-8')

    # 4. Embed fonts
    user_fonts = os.path.expanduser('~/Library/Fonts')
    font_files = {
        'Geist':         {'regular': os.path.join(user_fonts, 'Geist-Regular.otf'),
                          'bold':    os.path.join(user_fonts, 'Geist-Bold.otf')},
        'Geist ExtraBold': {'regular': os.path.join(user_fonts, 'Geist-ExtraBold.otf')},
        'DM Sans':       {'regular': os.path.join(user_fonts, 'DMSans[opsz,wght].ttf')},
    }

    def _obfuscate(font_data, guid_str):
        key  = bytes.fromhex(guid_str.strip('{}').replace('-', ''))
        data = bytearray(font_data)
        for i in range(min(32, len(data))):
            data[i] ^= key[i % 16]
        return bytes(data)

    rids      = [int(x) for x in re.findall(r'rId(\d+)', rels_xml)]
    next_rid  = max(rids) + 1 if rids else 30
    font_num  = 1
    new_rels  = []
    font_nodes = []

    for typeface, variants in font_files.items():
        var_parts = []
        for variant_key, fpath in variants.items():
            if not os.path.exists(fpath):
                print(f'  [embed] SKIP missing: {fpath}')
                continue
            with open(fpath, 'rb') as fh:
                raw = fh.read()
            g   = '{' + str(uuid.uuid4()).upper() + '}'
            obf = _obfuscate(raw, g)
            fname = f'ppt/fonts/font{font_num}.fntdata'
            files[fname] = obf
            rid = f'rId{next_rid}'
            new_rels.append(
                f'<Relationship Id="{rid}" '
                f'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
                f'relationships/font" Target="fonts/font{font_num}.fntdata"/>'
            )
            var_parts.append(f'<p:{variant_key} r:id="{rid}" uniqueId="{g}"/>')
            font_num += 1
            next_rid += 1
        if var_parts:
            font_nodes.append(
                f'<p:embeddedFont>'
                f'<p:font typeface="{typeface}" charset="0" pitchFamily="34"/>'
                + ''.join(var_parts) +
                f'</p:embeddedFont>'
            )
            print(f'  [embed] {typeface}: {list(variants.keys())}')

    if font_nodes:
        ct_xml = files['[Content_Types].xml'].decode('utf-8')
        if 'fntdata' not in ct_xml:
            ct_xml = ct_xml.replace(
                '</Types>',
                '<Default Extension="fntdata" ContentType="application/x-fontdata"/></Types>'
            )
            files['[Content_Types].xml'] = ct_xml.encode('utf-8')
        rels_xml = rels_xml.replace('</Relationships>',
                                    '\n'.join(new_rels) + '\n</Relationships>')
        embed_block = '<p:embeddedFontLst>' + ''.join(font_nodes) + '</p:embeddedFontLst>'
        if 'embeddedFontLst' not in prs_xml:
            prs_xml = prs_xml.replace('</p:presentation>', embed_block + '</p:presentation>')
    else:
        print('  [embed] no fonts embedded — check ~/Library/Fonts for Geist-Regular.otf etc.')

    files['ppt/presentation.xml']            = prs_xml.encode('utf-8')
    files['ppt/_rels/presentation.xml.rels'] = rels_xml.encode('utf-8')

    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for fname, data in files.items():
            zout.writestr(infos[fname] if fname in infos else fname, data)
